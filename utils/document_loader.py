"""
统一的文档加载器 —— 支持 PDF / DOCX / XLSX / PPTX / MD / TXT 多格式解析
返回 LangChain Document 对象列表

用法：
    from utils.document_loader import load_document
    docs = load_document("example.xlsx")
    text = docs[0].page_content
    source = docs[0].metadata["source"]
"""

import os
import re
import subprocess
import unicodedata
import logging
from typing import List

from langchain_core.documents import Document


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 统一入口
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def load_document(file_path: str) -> List[Document]:
    """
    统一文档加载入口，根据文件后缀自动分发到对应解析器

    Args:
        file_path: 文件绝对路径

    Returns:
        List[Document]: LangChain Document 对象列表

    Raises:
        FileNotFoundError: 文件不存在
        ValueError: 不支持的文件格式
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    file_name = os.path.basename(file_path)

    loaders = {
        ".pdf": _load_pdf,
        ".docx": _load_docx,
        ".xlsx": _load_excel,
        ".xls": _load_excel,
        ".pptx": _load_pptx,
        ".md": _load_markdown,
        ".markdown": _load_markdown,
        ".txt": _load_txt,
    }

    loader = loaders.get(ext)
    if loader is None:
        raise ValueError(f"不支持的文件格式: {ext}，支持的格式：{list(loaders.keys())}")

    logging.info(f"正在解析文档: {file_name}")
    docs = loader(file_path)

    # 确保每个 Document 都携带基本元数据
    for doc in docs:
        if "source" not in doc.metadata:
            doc.metadata["source"] = file_name
        if "file_path" not in doc.metadata:
            doc.metadata["file_path"] = file_path
        if "file_type" not in doc.metadata:
            doc.metadata["file_type"] = ext.lstrip(".")

    logging.info(f"文档解析完成: {file_name} → {len(docs)} 个 Document")
    return docs


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PDF —— pdfminer + OCR 后备（扫描版 PDF）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _load_pdf(file_path: str) -> List[Document]:
    """
    加载 PDF 文件。
    优先使用 pdfminer 提取文本；如果结果为空（扫描版 PDF），
    自动降级尝试 OCR 识别。
    """
    text = ""
    ocr_used = False

    # 1️⃣ 先用 pdfminer 提取文本
    try:
        from pdfminer.high_level import extract_text_to_fp
        from io import StringIO

        output = StringIO()
        with open(file_path, "rb") as f:
            extract_text_to_fp(f, output)
        text = output.getvalue()
    except Exception as e:
        logging.warning(f"pdfminer 解析异常: {e}")

    # 2️⃣ 如果结果为空或只有空白 → 尝试 OCR
    if not text.strip():
        logging.info(f"pdfminer 未提取到文本，尝试 OCR 识别: {os.path.basename(file_path)}")
        ocr_text = _ocr_pdf(file_path)
        if ocr_text.strip():
            # 检查 OCR 结果是否为乱码（字体损坏）
            if _is_garbled_text(ocr_text):
                logging.warning(f"OCR 结果疑似乱码（字体损坏）: {file_path}")
                raise ValueError("OCR 识别结果异常：PDF 字体可能已损坏或编码不兼容，请检查文件")
            text = ocr_text
            ocr_used = True
        else:
            logging.warning(f"PDF 文件为空且 OCR 也未能识别: {file_path}")
            return []

    metadata = {
        "source": os.path.basename(file_path),
        "file_path": file_path,
        "file_type": "pdf",
    }
    if ocr_used:
        metadata["ocr"] = True
        metadata["parser"] = "tesseract"

    return [Document(page_content=text, metadata=metadata)]


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# OCR 引擎（供 PDF / DOCX 内嵌图片共用）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

_OCR_AVAILABLE = None  # 延迟检测状态


def _check_ocr_available() -> bool:
    """
    检测 Tesseract OCR 引擎是否可用（仅检测一次）

    检测优先级（三级兜底）：
    1. TESSERACT_CMD 配置的绝对路径文件是否存在
    2. 系统命令行 `tesseract -v` 是否可用
    3. pytesseract.get_tesseract_version() 调用是否成功
    """
    global _OCR_AVAILABLE
    if _OCR_AVAILABLE is not None:
        return _OCR_AVAILABLE

    try:
        import pytesseract as ts

        # ── 第一级：配置的绝对路径是否存在 ──
        configured_cmd = getattr(ts.pytesseract, "tesseract_cmd", "") or ""
        # 也检查环境变量 TESSERACT_CMD（兼容 config.py 未提前设置的情况）
        env_cmd = os.environ.get("TESSERACT_CMD", "").strip()
        if env_cmd and os.path.isfile(env_cmd):
            configured_cmd = env_cmd
            ts.pytesseract.tesseract_cmd = env_cmd
        if configured_cmd and os.path.isfile(configured_cmd):
            _OCR_AVAILABLE = True
            logging.info(f"✅ Tesseract OCR 引擎可用（配置路径: {configured_cmd}）")
            return _OCR_AVAILABLE

        # ── 第二级：系统命令行 tesseract -v 兜底 ──
        try:
            result = subprocess.run(
                ["tesseract", "-v"],
                capture_output=True, text=True, timeout=10
            )
            if result.returncode == 0:
                # 找到系统 tesseract，设置到 pytesseract
                try:
                    import shutil
                    sys_path = shutil.which("tesseract")
                    if sys_path:
                        ts.pytesseract.tesseract_cmd = sys_path
                        logging.info(f"✅ Tesseract OCR 引擎可用（系统路径: {sys_path}）")
                except Exception:
                    pass
                _OCR_AVAILABLE = True
                return _OCR_AVAILABLE
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # ── 第三级：pytesseract 自带的检测 ──
        ts.get_tesseract_version()
        _OCR_AVAILABLE = True
        logging.info("✅ Tesseract OCR 引擎可用（pytesseract 自动检测）")

    except Exception:
        _OCR_AVAILABLE = False
        logging.warning(
            "⚠️ Tesseract OCR 引擎不可用，扫描 PDF / 图片将跳过 OCR。"
            " 请安装 Tesseract: https://github.com/tesseract-ocr/tesseract\n"
            "   或在 .env 中配置 TESSERACT_CMD 指向 tesseract.exe 的绝对路径"
        )

    return _OCR_AVAILABLE


def is_ocr_available() -> bool:
    """公开的 OCR 可用性检查（供上层 UI 调用）"""
    return _check_ocr_available()


_KNOWN_GARBLED_CHARS = set(
    "�￾￿"                          # 替换字符 / 非字符
    "□■▣▤▥▦"        # 方框类几何形状
    "△▲○●◇◆"        # 三角形 / 圆形 / 菱形
    "▽▼⬛⬜"                     # 更多几何形状
)


def _is_garbled_text(text: str, threshold: float = 0.3) -> bool:
    """
    检测 OCR 结果是否为大面积乱码（方框、替换字符等）。

    扫描 PDF 时若字体损坏 / 编码不兼容，Tesseract 可能输出
    大量不可读的占位字符，此时应判定解析失败而非返回乱码。

    Args:
        text: OCR 识别文本
        threshold: 乱码字符占可见字符的比例阈值

    Returns:
        True 表示文本疑似损坏/乱码
    """
    if not text.strip():
        return True

    visible_chars = [c for c in text if not c.isspace() and unicodedata.category(c) != "Cc"]
    if not visible_chars:
        return True

    garbled_in_visible = sum(1 for c in visible_chars if c in _KNOWN_GARBLED_CHARS)

    # 额外检查：正常 CJK / 拉丁字符占比过低也是乱码信号
    normal_chars = sum(
        1 for c in visible_chars
        if '一' <= c <= '鿿'          # CJK 统一表意文字
        or '぀' <= c <= 'ヿ'           # 日文假名
        or '가' <= c <= '힯'           # 韩文
        or 'a' <= c.lower() <= 'z'             # 拉丁字母
        or c.isdigit()                         # 数字
        or c in '.,;:!?。，；：！？、""''（）【】《》·—…'     # 常见标点
    )
    normal_ratio = normal_chars / len(visible_chars)

    return (garbled_in_visible / len(visible_chars)) > threshold or normal_ratio < 0.1


def _ocr_image_bytes(img_bytes: bytes, page_label: str = "") -> str:
    """对单张图片字节数据执行 OCR"""
    if not _check_ocr_available():
        return ""
    try:
        from PIL import Image
        import pytesseract
        from io import BytesIO

        img = Image.open(BytesIO(img_bytes))
        # 转 RGB（灰度图或 RGBA 可能影响识别）
        if img.mode != "RGB":
            img = img.convert("RGB")
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        text = text.strip()
        if text:
            prefix = f"[{page_label}] " if page_label else ""
            return f"{prefix}{text}"
        return ""
    except Exception as e:
        logging.warning(f"图片 OCR 失败: {e}")
        return ""


def _ocr_pdf(file_path: str) -> str:
    """
    对扫描版 PDF 执行 OCR（整套图片 → 文字）
    需要系统安装 Tesseract 引擎 + Python 包 pdf2image + pytesseract
    """
    if not _check_ocr_available():
        return ""

    try:
        from pdf2image import convert_from_path

        images = convert_from_path(file_path, dpi=300)
        if not images:
            return ""

        import pytesseract as ts

        all_text = []
        for i, img in enumerate(images):
            try:
                if img.mode != "RGB":
                    img = img.convert("RGB")
                page_text = ts.image_to_string(img, lang="chi_sim+eng")
                page_text = page_text.strip()
                if page_text:
                    all_text.append(f"--- OCR 第 {i + 1} 页 ---\n{page_text}")
            except Exception as e:
                logging.warning(f"OCR 第 {i + 1} 页失败: {e}")

        return "\n\n".join(all_text)

    except ImportError as e:
        pkg = "pdf2image" if "pdf2image" in str(e) else "pytesseract"
        logging.warning(f"OCR 需要安装 {pkg}: pip install {pkg}")
        return ""
    except Exception as e:
        err_msg = str(e).lower()
        if "tesseract" in err_msg or "poppler" in err_msg:
            logging.warning(
                "⚠️ OCR 引擎组件缺失。Tesseract: https://github.com/tesseract-ocr/tesseract | "
                "Poppler: https://poppler.freedesktop.org"
            )
        else:
            logging.error(f"PDF OCR 识别失败: {e}")
        return ""


def _ocr_images_from_docx(file_path: str) -> List[str]:
    """
    从 docx 压缩包中提取内嵌图片并 OCR
    docx 本质是 ZIP，图片通常存放在 word/media/ 下
    """
    if not _check_ocr_available():
        return []

    import zipfile
    from io import BytesIO

    results = []
    try:
        with zipfile.ZipFile(file_path, "r") as z:
            image_names = [
                n for n in z.namelist()
                if n.startswith("word/media/") and n.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".gif"))
            ]
            if not image_names:
                return []

            for name in sorted(image_names):
                img_bytes = z.read(name)
                text = _ocr_image_bytes(img_bytes, page_label=f"图片: {os.path.basename(name)}")
                if text:
                    results.append(text)

        return results
    except zipfile.BadZipFile:
        logging.warning(f"DOCX 文件损坏，无法提取图片: {file_path}")
        return []
    except Exception as e:
        logging.warning(f"DOCX 图片提取失败: {e}")
        return []


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# DOCX —— 段落 + 表格 + 内嵌图片 OCR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _load_docx(file_path: str) -> List[Document]:
    """加载 DOCX 文件：段落 + 表格 + 内嵌图片 OCR"""
    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)

        # 提取所有段落文本
        paragraphs = [p.text for p in doc.paragraphs]

        # 提取表格内容
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                tables_text.append(" | ".join(cells))

        all_text = "\n".join(paragraphs)
        if tables_text:
            all_text += "\n\n【表格内容】\n" + "\n".join(tables_text)

        # 提取内嵌图片 OCR 文字
        ocr_texts = _ocr_images_from_docx(file_path)
        if ocr_texts:
            all_text += "\n\n【图片OCR识别】\n" + "\n".join(ocr_texts)

        if not all_text.strip():
            logging.warning(f"DOCX 文件内容为空: {file_path}")
            return []

        return [
            Document(
                page_content=all_text,
                metadata={
                    "source": os.path.basename(file_path),
                    "file_path": file_path,
                    "file_type": "docx",
                    "paragraphs": len(paragraphs),
                    "tables": len(tables_text),
                    "ocr_images": len(ocr_texts),
                },
            )
        ]
    except ImportError:
        logging.error("处理 DOCX 文件需要安装 python-docx 库")
        raise
    except Exception as e:
        logging.error(f"DOCX 解析失败: {file_path}, 错误: {e}")
        raise


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# XLSX / XLS —— openpyxl 逐工作表、逐单元格读取
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _load_excel(file_path: str) -> List[Document]:
    """
    加载 Excel 文件（.xlsx / .xls）。
    遍历所有工作表，逐单元格读取内容，拼接为完整文本。
    """
    try:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        all_sheets_text = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                # 过滤掉全空行
                cells = [str(c).strip() for c in row if c is not None]
                if cells:
                    rows_text.append(" | ".join(cells))

            if rows_text:
                sheet_content = "\n".join(rows_text)
                all_sheets_text.append(f"【工作表: {sheet_name}】\n{sheet_content}")

        wb.close()

        if not all_sheets_text:
            logging.warning(f"Excel 文件内容为空: {file_path}")
            return []

        # 如果只有一个工作表，返回单个 Document；多个则每个工作表返回一个 Document
        if len(all_sheets_text) == 1:
            return [
                Document(
                    page_content=all_sheets_text[0],
                    metadata={
                        "source": os.path.basename(file_path),
                        "file_path": file_path,
                        "file_type": "excel",
                        "sheets": 1,
                    },
                )
            ]
        else:
            docs = []
            for sheet_text in all_sheets_text:
                sheet_name = sheet_text.split("\n")[0].replace("【工作表: ", "").replace("】", "")
                docs.append(
                    Document(
                        page_content=sheet_text,
                        metadata={
                            "source": os.path.basename(file_path),
                            "file_path": file_path,
                            "file_type": "excel",
                            "sheet": sheet_name,
                        },
                    )
                )
            return docs

    except ImportError:
        logging.error("处理 Excel 文件需要安装 openpyxl 库")
        raise
    except Exception as e:
        logging.error(f"Excel 解析失败: {file_path}, 错误: {e}")
        raise


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PPTX —— python-pptx 提取每页文本框文字
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _load_pptx(file_path: str) -> List[Document]:
    """
    加载 PPTX 文件。
    遍历每张幻灯片，提取所有形状（文本框）内的文字内容。
    """
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            page_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            page_texts.append(para_text)

            if page_texts:
                slide_content = "\n".join(page_texts)
                slides_text.append(f"--- 第 {slide_idx} 页 ---\n{slide_content}")

        if not slides_text:
            logging.warning(f"PPT 文件内容为空: {file_path}")
            return []

        combined = "\n\n".join(slides_text)

        return [
            Document(
                page_content=combined,
                metadata={
                    "source": os.path.basename(file_path),
                    "file_path": file_path,
                    "file_type": "pptx",
                    "slides": len(slides_text),
                },
            )
        ]
    except ImportError:
        logging.error("处理 PPT 文件需要安装 python-pptx 库")
        raise
    except Exception as e:
        logging.error(f"PPT 解析失败: {file_path}, 错误: {e}")
        raise


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MD / TXT —— 文本类（沿用原有实现）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _load_markdown(file_path: str) -> List[Document]:
    """加载 Markdown 文件（直接读取 UTF-8 文本）"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

        if not text.strip():
            logging.warning(f"Markdown 文件内容为空: {file_path}")
            return []

        return [
            Document(
                page_content=text,
                metadata={
                    "source": os.path.basename(file_path),
                    "file_path": file_path,
                    "file_type": "markdown",
                },
            )
        ]
    except Exception as e:
        logging.error(f"Markdown 读取失败: {file_path}, 错误: {e}")
        raise


def _load_txt(file_path: str) -> List[Document]:
    """加载 TXT 文件"""
    try:
        from langchain_community.document_loaders import TextLoader

        loader = TextLoader(file_path, encoding="utf-8")
        docs = loader.load()

        for doc in docs:
            doc.metadata.setdefault("source", os.path.basename(file_path))
            doc.metadata.setdefault("file_path", file_path)
            doc.metadata.setdefault("file_type", "txt")

        return docs
    except ImportError:
        logging.warning("langchain_community 未安装，使用内置方式读取 TXT")
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()

            if not text.strip():
                return []

            return [
                Document(
                    page_content=text,
                    metadata={
                        "source": os.path.basename(file_path),
                        "file_path": file_path,
                        "file_type": "txt",
                    },
                )
            ]
        except Exception as e:
            logging.error(f"TXT 读取失败: {file_path}, 错误: {e}")
            raise
    except Exception as e:
        logging.error(f"TXT 读取失败: {file_path}, 错误: {e}")
        raise
