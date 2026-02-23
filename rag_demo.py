import os
os.environ['HF_ENDPOINT'] = 'https://hf-mirror.com'  # 设置 Hugging Face 镜像源，解决国内无法直接下载模型的问题
import gradio as gr  # 构建 Web UI 的库
from pdfminer.high_level import extract_text_to_fp  # 从 PDF 中提取文本
from sentence_transformers import SentenceTransformer  # 文本向量化模型
from sentence_transformers import CrossEncoder  # 导入交叉编码器，用于细粒度语义匹配（如重排序阶段）
from faiss import IndexFlatL2, IndexIVFFlat, IndexIVFPQ # Facebook AI 的向量相似度检索库（用于构建向量索引）
import requests  # 网络请求库，用于 API 请求、模型下载等
import json  # 处理 JSON 格式数据
from io import StringIO  # 将 PDF 提取内容写入内存对象中
from langchain_text_splitters import RecursiveCharacterTextSplitter  # 文本分段工具
import os  # 再次导入 os（可省略，属于冗余）
import socket  # 判断端口占用或本地网络检查
import webbrowser  # 自动打开 Web UI 页面
import logging  # 日志记录，用于调试与错误跟踪
from requests.adapters import HTTPAdapter  # 请求重试机制
from urllib3.util.retry import Retry
import time  # 时间处理，用于计时或延迟
from datetime import datetime  # 获取时间戳、日志时间等
import re  # 正则处理，适用于 PDF 内容清洗等
from dotenv import load_dotenv  # 加载环境变量配置（如 API KEY）
from rank_bm25 import BM25Okapi  # 导入 BM25 算法，用于传统稀疏检索，与向量搜索可混合增强召回率
import numpy as np  # 向量计算、FAISS 依赖的数值库
import jieba  # 中文分词库，用于 BM25 中文检索效果提升
import threading  # 多线程加速处理，例如文档向量化
from pathlib import Path  # 高级文件路径处理
from functools import lru_cache  # 缓存装饰器，提高重复调用性能
from typing import List, Tuple, Any, Optional  # 类型注解，提高代码可读性与维护性


# 加载环境变量
# 加载 .env 配置文件路径
dotenv_path = Path(__file__).parent / "example.env"
load_dotenv(dotenv_path)
# 从 .env 文件中读取 SERPAPI_KEY 用于搜索引擎查询
SERPAPI_KEY = os.getenv("SERPAPI_KEY")  # 在 .env 文件中设置 SERPAPI_KEY
# 默认搜索引擎为 Google，如需可切换至其他搜索源
SEARCH_ENGINE = "google"
RERANK_METHOD = os.getenv("RERANK_METHOD", "cross_encoder")

#  SiliconFlow API 配置，用于 LLM 调用
SILICONFLOW_API_KEY = os.getenv("SILICONFLOW_API_KEY")  # 在 .env 设置 API Key
SILICONFLOW_API_URL = os.getenv(
    "SILICONFLOW_API_URL",
    "https://api.siliconflow.cn/v1/chat/completions"
)  # 默认访问 SiliconFlow Chat API

# 模型名称配置，支持通过环境变量自定义本地和云端模型
OLLAMA_MODEL_NAME = os.getenv("OLLAMA_MODEL_NAME", "deepseek-ai/DeepSeek-R1-Distill-Qwen-7B")
SILICONFLOW_MODEL_NAME = os.getenv("SILICONFLOW_MODEL_NAME", "deepseek-ai/DeepSeek-R1-Distill-Qwen-7B")

# 设置请求超时与重试机制，提升网络访问的稳定性
requests.adapters.DEFAULT_RETRIES = 3  # 增加网络请求失败后的重试次数

# 设置环境变量，关闭 TensorFlow 的 oneDNN 优化，避免某些 CPU 环境下的兼容性问题
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'

# 配置代理绕过规则，避免本地访问（如 127.0.0.1）走代理导致无法连接
os.environ['NO_PROXY'] = 'localhost,127.0.0.1'

# 初始化向量嵌入模型，用于将文本转换为向量(英文优化模型)
EMBED_MODEL = SentenceTransformer('all-MiniLM-L6-v2')
# 若主要处理中文文档，可切换为中文优化模型，例如：
# EMBED_MODEL = SentenceTransformer('shibing624/text2vec-base-chinese')

# FAISS 相关的全局数据结构
faiss_index = None  # FAISS 索引对象，存储向量用于相似度检索
faiss_contents_map = {}  # original_id -> content 映射，用于检索结果的原文还原
faiss_metadatas_map = {}  # original_id -> metadata 映射，用于记录文档片段元信息
faiss_id_order_for_index = []  # 记录向量添加到 FAISS 的顺序，便于管理和增量更新

# 初始化交叉编码器，用于重排序结果。采取懒加载策略，首次使用时再加载模型
cross_encoder = None
cross_encoder_lock = threading.Lock()  # 线程锁，避免多线程环境中重复加载交叉编码器

def get_cross_encoder():
    """懒加载交叉编码器模型，避免首次运行阻塞"""
    global cross_encoder
    if cross_encoder is None:
        with cross_encoder_lock:
            if cross_encoder is None:
                try:
                    # 多语言交叉编码器，更适合中文语义重排序,专门用来对「初步检索回来的候选结果」做更精细的语义相关性打分，从而提升最终排序质量
                    cross_encoder = CrossEncoder(
                        'sentence-transformers/distiluse-base-multilingual-cased-v2'
                    )
                    logging.info("交叉编码器加载成功")
                except Exception as e:
                    logging.error(f"加载交叉编码器失败: {str(e)}")
                    # 保持 None，以便后续重试加载
                    cross_encoder = None
    return cross_encoder

# 新增：自动选择FAISS索引类型的封装类
class AutoFaissIndex:
    def __init__(self, dimension=384):
        """
        自动选择FAISS索引类型的封装类
        参数:
            dimension: 向量维度 (默认为384，对应all-MiniLM-L6-v2模型的输出维度)
        """
        self.dimension = dimension
        self.index = None
        self.index_type = None
        self.nlist = None  # IVF类索引的聚类中心数
        self.m = None  # PQ类索引的细分段数
        self.nprobe = None  # IVF类索引的搜索聚类中心数

        # 阈值配置 (可根据实际硬件调整)
        self.small_dataset_threshold = 10_000  # 小数据集阈值
        self.medium_dataset_threshold = 100_000  # 中等数据集阈值
        self.large_dataset_threshold = 1_000_000  # 大数据集阈值

    @property
    def ntotal(self):
        """返回索引中的向量总数"""
        return self.index.ntotal if self.index else 0

    def select_index_type(self, num_vectors):
        """
        根据向量数量自动选择最优索引类型
        参数:
            num_vectors: 要索引的向量数量
        """
        if num_vectors <= self.small_dataset_threshold:
            # 小数据集: 使用精确搜索的Flat索引
            self.index_type = "FlatL2"
            self.index = IndexFlatL2(self.dimension)
            self.nprobe = 1  # 不相关，仅为统一接口

        elif num_vectors <= self.medium_dataset_threshold:
            # 中等数据集: 使用IVFFlat平衡精度和速度
            self.index_type = "IVFFlat"
            self.nlist = min(100, int(np.sqrt(num_vectors)))
            quantizer = IndexFlatL2(self.dimension)
            self.index = IndexIVFFlat(quantizer, self.dimension, self.nlist)
            self.nprobe = min(10, max(1, int(self.nlist * 0.1)))  # 搜索10%的聚类中心

        else:
            # 大数据集: 使用IVFPQ牺牲少量精度换取更高效率
            self.index_type = "IVFPQ"
            self.nlist = min(256, int(np.sqrt(num_vectors)))
            self.m = min(8, self.dimension // 4)  # 每个向量分成8段
            quantizer = IndexFlatL2(self.dimension)
            self.index = IndexIVFPQ(quantizer, self.dimension, self.nlist, self.m, 8)  # 8 bits per code
            self.nprobe = min(32, max(1, int(self.nlist * 0.05)))  # 搜索5%的聚类中心

        return self.index_type

    def train(self, vectors):
        """
        训练索引 (仅IVF类索引需要)

        参数:
            vectors: 用于训练的向量数组 (np.array)
        """
        if self.index_type in ["IVFFlat", "IVFPQ"]:
            self.index.train(vectors)

    def add(self, vectors):
        """
        添加向量到索引

        参数:
            vectors: 要添加的向量数组 (np.array)
        """
        if self.index_type in ["IVFFlat", "IVFPQ"] and not self.index.is_trained:
            self.train(vectors)

        self.index.add(vectors)

    def search(self, query_vectors, k=5):
        """
        执行搜索

        参数:
            query_vectors: 查询向量数组 (np.array)
            k: 返回的最近邻数量

        返回:
            distances: 距离矩阵 (nq, k)
            indices: 索引矩阵 (nq, k)
        """
        if self.index_type in ["IVFFlat", "IVFPQ"]:
            self.index.nprobe = self.nprobe

        return self.index.search(query_vectors, k)

    def get_index_info(self):
        """获取当前索引配置信息"""
        return {
            "index_type": self.index_type,
            "dimension": self.dimension,
            "nlist": self.nlist,
            "m": self.m,
            "nprobe": self.nprobe,
            "size": self.ntotal
        }

# 递归检索主逻辑
# 多轮检索 → 混合排序 → 交叉编码器 → LLM 引导 Query 扩展
def recursive_retrieval(initial_query, max_iterations=3, enable_web_search=False, model_choice="siliconflow"):
    """
    递归检索与查询迭代优化
    利用现有检索结果，通过 LLM 判断是否继续改写查询进行下一轮检索

    Args:
        initial_query: 初始查询文本
        max_iterations: 最大检索迭代轮数
        enable_web_search: 是否启用互联网补充搜索结果
        model_choice: 使用推理模型来源("ollama" 或 "siliconflow(硅基流动)")

    Returns:
        all_contexts: 所有获取到的文本内容列表
        all_doc_ids: 文本对应的原始 ID 列表
        all_metadata: 文本元信息列表
    """
    query = initial_query
    all_contexts = []
    all_doc_ids = []
    all_metadata = []

    global faiss_index, faiss_contents_map, faiss_metadatas_map, faiss_id_order_for_index

    for i in range(max_iterations):
        logging.info(f"递归检索 {i + 1}/{max_iterations}，当前 Query: {query}")

        # 存储来自网络搜索的补充信息，优先提供给 LLM 判断
        web_results_texts = []

        if enable_web_search and check_serpapi_key():
            try:
                web_search_raw_results = update_web_results(query)
                for res in web_search_raw_results:
                    # 仅加入内容，不加入向量库，避免在线来源污染本地索引
                    text = f"标题：{res.get('title', '')}\n摘要：{res.get('snippet', '')}"
                    web_results_texts.append(text)
            except Exception as e:
                logging.error(f"网络搜索出错: {str(e)}")

        # 语义检索
        query_embedding = EMBED_MODEL.encode([query])
        query_embedding_np = np.array(query_embedding).astype('float32')

        semantic_results_docs = []
        semantic_results_metadatas = []
        semantic_results_ids = []

        # 修复FAISS检索部分
        if faiss_index is not None and hasattr(faiss_index, 'ntotal') and faiss_index.ntotal > 0:
            try:
                D, I = faiss_index.search(query_embedding_np, k=10)
                # 将FAISS索引转换回原始ID
                for faiss_idx in I[0]:
                    if faiss_idx != -1 and faiss_idx < len(faiss_id_order_for_index):
                        original_id = faiss_id_order_for_index[faiss_idx]
                        if original_id in faiss_contents_map:  # 添加存在性检查
                            semantic_results_docs.append(faiss_contents_map.get(original_id, ""))
                            semantic_results_metadatas.append(faiss_metadatas_map.get(original_id, {}))
                            semantic_results_ids.append(original_id)
                        else:
                            logging.warning(f"ID {original_id} 不在内容映射中")
            except Exception as e:
                logging.error(f"FAISS 检索错误: {str(e)}")
        else:
            logging.warning("FAISS索引为空或未初始化")

        # 稀疏检索 BM25
        bm25_results = BM25_MANAGER.search(query, top_k=10) if BM25_MANAGER.bm25_index else []

        # 格式对齐，适配混合排序函数
        prepared_semantic_results_for_hybrid = {
            "ids": [semantic_results_ids],
            "documents": [semantic_results_docs],
            "metadatas": [semantic_results_metadatas]
        }

        # 混合语义 + BM25 排序
        hybrid_results = hybrid_merge(prepared_semantic_results_for_hybrid, bm25_results, alpha=0.7)

        doc_ids_current_iter = []
        docs_current_iter = []
        metadata_list_current_iter = []

        # 取 Top-10 进入交叉编码重排序
        if hybrid_results:
            for doc_id, result_data in hybrid_results[:10]:
                doc_ids_current_iter.append(doc_id)
                docs_current_iter.append(result_data['content'])
                metadata_list_current_iter.append(result_data['metadata'])

        # 交叉编码器重排序，提升准确性
        if docs_current_iter:
            try:
                reranked_results = rerank_results(query, docs_current_iter, doc_ids_current_iter,
                                                  metadata_list_current_iter, top_k=5)
            except Exception as e:
                logging.error(f"重排序失败: {str(e)}")
                # 回退为混合排序结果
                reranked_results = [
                    (doc_id, {'content': doc, 'metadata': meta, 'score': 1.0})
                    for doc_id, doc, meta in zip(doc_ids_current_iter, docs_current_iter, metadata_list_current_iter)
                ]
        else:
            reranked_results = []

        # 整合本轮检索结果
        current_contexts_for_llm = web_results_texts[:]
        for doc_id, result_data in reranked_results:
            if doc_id not in all_doc_ids:
                all_doc_ids.append(doc_id)
                all_contexts.append(result_data['content'])
                all_metadata.append(result_data['metadata'])
            current_contexts_for_llm.append(result_data['content'])

        if i == max_iterations - 1:
            break

        # 调用 LLM 决策是否生成新的查询
        if current_contexts_for_llm:
            current_summary = "\n".join(current_contexts_for_llm[:3])

            next_query_prompt = f"""你是一个查询优化助手。根据以下信息判断是否需要新的查询。

[初始问题]
{initial_query}

[检索结果摘要]
{current_summary}

要求：
1. 如果信息已足够，直接回复：不需要进一步查询
2. 否则返回一个更精准的新查询，仅包含查询词
"""

            try:
                if model_choice == "siliconflow":
                    logging.info("使用 SiliconFlow API 分析下一步查询")
                    result = call_siliconflow_api(next_query_prompt)
                    next_query = result.strip() if isinstance(result, str) else result[0].strip()

                    if "<think>" in next_query:
                        next_query = next_query.split("<think>")[0].strip()

                else:
                    logging.info("使用本地 Ollama 模型分析下一步查询")
                    response = session.post(
                        "http://localhost:11434/api/generate",
                        json={
                            "model": OLLAMA_MODEL_NAME,
                            "prompt": next_query_prompt,
                            "stream": False
                        },
                        timeout=180
                    )
                    next_query = response.json().get("response", "").strip()

                if "不需要" in next_query:
                    logging.info("LLM 判断无需更多查询")
                    break

                if len(next_query) > 100:
                    logging.warning("生成内容过长，不视为有效查询")
                    break

                query = next_query
                logging.info(f"生成下一轮查询: {query}")

            except Exception as e:
                logging.error(f"生成新查询失败: {str(e)}")
                break
        else:
            break

    return all_contexts, all_doc_ids, all_metadata


class BM25IndexManager:
    def __init__(self):
        self.bm25_index = None
        self.doc_mapping = {}  # 映射BM25索引位置到文档ID
        self.tokenized_corpus = []
        self.raw_corpus = []

    def build_index(self, documents, doc_ids):
        """构建BM25索引"""
        self.raw_corpus = documents
        self.doc_mapping = {i: doc_id for i, doc_id in enumerate(doc_ids)}

        # 对文档进行分词，使用jieba分词器更适合中文
        self.tokenized_corpus = []
        for doc in documents:
            # 对中文文档进行分词
            tokens = list(jieba.cut(doc))
            self.tokenized_corpus.append(tokens)

        # 创建BM25索引
        self.bm25_index = BM25Okapi(self.tokenized_corpus)
        return True

    def search(self, query, top_k=5):
        """使用BM25检索相关文档"""
        if not self.bm25_index:
            return []

        # 对查询进行分词
        tokenized_query = list(jieba.cut(query))

        # 获取BM25得分
        bm25_scores = self.bm25_index.get_scores(tokenized_query)

        # 获取得分最高的文档索引
        top_indices = np.argsort(bm25_scores)[-top_k:][::-1]

        # 返回结果
        results = []
        for idx in top_indices:
            if bm25_scores[idx] > 0:  # 只返回有相关性的结果
                results.append({
                    'id': self.doc_mapping[idx],
                    'score': float(bm25_scores[idx]),
                    'content': self.raw_corpus[idx]
                })

        return results

    def clear(self):
        """清空索引"""
        self.bm25_index = None
        self.doc_mapping = {}
        self.tokenized_corpus = []
        self.raw_corpus = []


# 初始化BM25索引管理器
BM25_MANAGER = BM25IndexManager()

logging.basicConfig(level=logging.INFO)

print("Gradio version:", gr.__version__)  # 添加版本输出

# 在初始化组件后添加：
session = requests.Session()
retries = Retry(
    total=3,
    backoff_factor=0.1,
    status_forcelist=[500, 502, 503, 504]
)
session.mount('http://', HTTPAdapter(max_retries=retries))


#########################################
# SerpAPI 网络查询及向量化处理函数
#########################################
def serpapi_search(query: str, num_results: int = 5) -> list:
    """
    执行 SerpAPI 搜索，并返回解析后的结构化结果
    """
    if not SERPAPI_KEY:
        raise ValueError("未设置 SERPAPI_KEY 环境变量。请在.env文件中设置您的 API 密钥。")
    try:
        params = {
            "engine": SEARCH_ENGINE,
            "q": query,
            "api_key": SERPAPI_KEY,
            "num": num_results,
            "hl": "zh-CN",  # 中文界面
            "gl": "cn"
        }
        response = requests.get("https://serpapi.com/search", params=params, timeout=15)
        response.raise_for_status()
        search_data = response.json()
        return _parse_serpapi_results(search_data)
    except Exception as e:
        logging.error(f"网络搜索失败: {str(e)}")
        return []


def _parse_serpapi_results(data: dict) -> list:
    """解析 SerpAPI 返回的原始数据"""
    results = []
    if "organic_results" in data:
        for item in data["organic_results"]:
            result = {
                "title": item.get("title"),
                "url": item.get("link"),
                "snippet": item.get("snippet"),
                "timestamp": item.get("date")  # 若有时间信息，可选
            }
            results.append(result)
    # 如果有知识图谱信息，也可以添加置顶（可选）
    if "knowledge_graph" in data:
        kg = data["knowledge_graph"]
        results.insert(0, {
            "title": kg.get("title"),
            "url": kg.get("source", {}).get("link", ""),
            "snippet": kg.get("description"),
            "source": "knowledge_graph"
        })
    return results


def update_web_results(query: str, num_results: int = 5) -> list:
    """
    基于 SerpAPI 搜索结果。注意：此版本不将结果存入FAISS。
    它仅返回原始搜索结果。
    """
    results = serpapi_search(query, num_results)
    if not results:
        logging.info("网络搜索没有返回结果或发生错误")
        return []

    # 之前这里有删除旧网络结果和添加到ChromaDB的逻辑。
    # 由于FAISS IndexFlatL2不支持按ID删除，并且动态添加涉及复杂ID管理，
    # 此简化版本不将网络结果添加到FAISS索引。
    # 返回原始结果，供调用者决定如何使用（例如，仅作为文本上下文）。
    logging.info(f"网络搜索返回 {len(results)} 条结果，这些结果不会被添加到FAISS索引中。")
    return results  # 返回原始SerpAPI结果列表


# 检查是否配置了SERPAPI_KEY
def check_serpapi_key():
    """检查是否配置了SERPAPI_KEY"""
    return SERPAPI_KEY is not None and SERPAPI_KEY.strip() != ""


# 添加文件处理状态跟踪
class FileProcessor:
    def __init__(self):
        self.processed_files = {}  # 存储已处理文件的状态

    def clear_files(self):
        """清空所有文件记录"""
        self.processed_files = {}

    def add_file(self, file_name):
        self.processed_files[file_name] = {
            'status': '等待处理',
            'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'chunks': 0
        }

    def update_status(self, file_name, status, chunks=None):
        if file_name in self.processed_files:
            self.processed_files[file_name]['status'] = status
            if chunks is not None:
                self.processed_files[file_name]['chunks'] = chunks

    def get_file_list(self):
        return [
            f"📄 {fname} | {info['status']}"
            for fname, info in self.processed_files.items()
        ]


file_processor = FileProcessor()


#########################################
# 矛盾检测函数
#########################################
def detect_conflicts(sources):
    """精准矛盾检测算法"""
    key_facts = {}
    for item in sources:
        facts = extract_facts(item['text'] if 'text' in item else item.get('excerpt', ''))
        for fact, value in facts.items():
            if fact in key_facts:
                if key_facts[fact] != value:
                    return True
            else:
                key_facts[fact] = value
    return False


def extract_facts(text):
    """从文本提取关键事实（示例逻辑）"""
    facts = {}
    # 提取数值型事实
    numbers = re.findall(r'\b\d{4}年|\b\d+%', text)
    if numbers:
        facts['关键数值'] = numbers
    # 提取技术术语
    if "产业图谱" in text:
        facts['技术方法'] = list(set(re.findall(r'[A-Za-z]+模型|[A-Z]{2,}算法', text)))
    return facts


def evaluate_source_credibility(source):
    """评估来源可信度"""
    credibility_scores = {
        "gov.cn": 0.9,
        "edu.cn": 0.85,
        "weixin": 0.7,
        "zhihu": 0.6,
        "baidu": 0.5
    }

    url = source.get('url', '')
    if not url:
        return 0.5  # 默认中等可信度

    domain_match = re.search(r'//([^/]+)', url)
    if not domain_match:
        return 0.5

    domain = domain_match.group(1)

    # 检查是否匹配任何已知域名
    for known_domain, score in credibility_scores.items():
        if known_domain in domain:
            return score

    return 0.5  # 默认中等可信度


# 修改后的 extract_text 函数, 支持多种文件格式
def extract_text(filepath):
    """支持多种文件格式的文本提取"""
    file_ext = os.path.splitext(filepath)[1].lower()

    if file_ext == '.pdf':
        output = StringIO()
        with open(filepath, 'rb') as file:
            extract_text_to_fp(file, output)
        return output.getvalue()

    elif file_ext in ['.txt', '.md']:
        with open(filepath, 'r', encoding='utf-8') as file:
            return file.read()

    elif file_ext in ['.docx']:
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        except ImportError:
            logging.error("处理Word文档需要安装python-docx库")
            return ""

    elif file_ext in ['.xlsx', '.xls']:
        try:
            import pandas as pd
            text = ""
            xl = pd.ExcelFile(filepath)
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                text += f"工作表: {sheet_name}\n"
                text += df.to_string(index=False) + "\n\n"
            return text
        except ImportError:
            logging.error("处理Excel文件需要安装pandas库")
            return ""

    elif file_ext in ['.pptx']:
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            logging.error("处理PPT文件需要安装python-pptx库")
            return ""

    else:
        logging.warning(f"不支持的文件格式: {file_ext}")
        return ""

# 新增：可以扩展多种文件格式
def process_multiple_pdfs(files: List[Any], progress=gr.Progress()):
    """处理多个文件"""
    if not files:
        return "请选择要上传的文件(支持PDF, Word, Excel, PPT, TXT, Markdown等)", []

    try:
        # 清空向量数据库和相关存储
        progress(0.1, desc="清理历史数据...")
        global faiss_index, faiss_contents_map, faiss_metadatas_map, faiss_id_order_for_index
        faiss_index = None
        faiss_contents_map = {}
        faiss_metadatas_map = {}
        faiss_id_order_for_index = []

        # 清空BM25索引
        BM25_MANAGER.clear()
        logging.info("成功清理历史FAISS数据和BM25索引")

        total_files = len(files)
        processed_results = []
        all_chunks = []
        all_metadatas = []
        all_ids = []

        for idx, file in enumerate(files, 1):
            try:
                file_name = os.path.basename(file.name)
                progress((idx - 1) / total_files, desc=f"处理文件 {idx}/{total_files}: {file_name}")

                # 提取文本
                text = extract_text(file.name)
                if not text:
                    raise ValueError("文档内容为空或无法提取文本")

                # 分块处理
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=400,
                    chunk_overlap=40,
                    separators=["\n\n", "\n", "。", "，", "；", "：", " ", ""]
                )
                chunks = text_splitter.split_text(text)

                # 生成唯一ID和元数据
                doc_id = f"doc_{int(time.time())}_{idx}"
                metadatas = [{"source": file_name, "doc_id": doc_id} for _ in chunks]
                chunk_ids = [f"{doc_id}_chunk_{i}" for i in range(len(chunks))]

                # 保存数据
                all_chunks.extend(chunks)
                all_metadatas.extend(metadatas)
                all_ids.extend(chunk_ids)

                processed_results.append(f"✅ {file_name}: 成功处理 {len(chunks)} 个文本块")

            except Exception as e:
                error_msg = str(e)
                logging.error(f"处理文件 {file_name} 时出错: {error_msg}")
                processed_results.append(f"❌ {file_name}: 处理失败 - {error_msg}")

        # 批量生成嵌入向量
        if all_chunks:
            progress(0.8, desc="生成文本嵌入...")
            embeddings = EMBED_MODEL.encode(all_chunks, show_progress_bar=True)
            embeddings_np = np.array(embeddings).astype('float32')

            # 构建FAISS索引
            progress(0.9, desc="构建FAISS索引...")
            dimension = embeddings_np.shape[1]
            faiss_index = IndexFlatL2(dimension)  # 使用基础索引

            # 确保内容映射同步
            for chunk_id, chunk, meta in zip(all_ids, all_chunks, all_metadatas):
                faiss_contents_map[chunk_id] = chunk
                faiss_metadatas_map[chunk_id] = meta
                faiss_id_order_for_index.append(chunk_id)

            faiss_index.add(embeddings_np)
            logging.info(f"FAISS索引构建完成，共索引 {faiss_index.ntotal} 个文本块")

        # 构建BM25索引
        progress(0.95, desc="构建BM25检索索引...")
        BM25_MANAGER.build_index(all_chunks, all_ids)

        summary = f"\n总计处理 {total_files} 个文件，{len(all_chunks)} 个文本块"
        processed_results.append(summary)

        return "\n".join(processed_results), [f"📄 {os.path.basename(f.name)}" for f in files]

    except Exception as e:
        error_msg = str(e)
        logging.error(f"整体处理过程出错: {error_msg}")
        return f"处理过程出错: {error_msg}", []

# 交叉编码器重排序函数（二次排序）
def rerank_with_cross_encoder(query, docs, doc_ids, metadata_list, top_k=5):
    """
    使用交叉编码器对检索结果进行重排序

    参数:
        query: 查询字符串
        docs: 文档内容列表
        doc_ids: 文档ID列表
        metadata_list: 元数据列表
        top_k: 返回结果数量

    返回:
        重排序后的结果列表 [(doc_id, {'content': doc, 'metadata': metadata, 'score': score}), ...]
    """
    if not docs:
        return []

    encoder = get_cross_encoder()
    if encoder is None:
        logging.warning("交叉编码器不可用，跳过重排序")
        # 返回原始顺序（按索引排序）
        return [(doc_id, {'content': doc, 'metadata': meta, 'score': 1.0 - idx / len(docs)})
                for idx, (doc_id, doc, meta) in enumerate(zip(doc_ids, docs, metadata_list))]

    # 准备交叉编码器输入
    cross_inputs = [[query, doc] for doc in docs]

    try:
        # 计算相关性得分
        scores = encoder.predict(cross_inputs)

        # 组合结果
        results = [
            (doc_id, {
                'content': doc,
                'metadata': meta,
                'score': float(score)  # 确保是Python原生类型
            })
            for doc_id, doc, meta, score in zip(doc_ids, docs, metadata_list, scores)
        ]

        # 按得分排序
        results = sorted(results, key=lambda x: x[1]['score'], reverse=True)

        # 返回前K个结果
        return results[:top_k]
    except Exception as e:
        logging.error(f"交叉编码器重排序失败: {str(e)}")
        # 出错时返回原始顺序
        return [(doc_id, {'content': doc, 'metadata': meta, 'score': 1.0 - idx / len(docs)})
                for idx, (doc_id, doc, meta) in enumerate(zip(doc_ids, docs, metadata_list))]


# LLM相关性评分函数
@lru_cache(maxsize=32)
def get_llm_relevance_score(query, doc):
    """
    使用LLM对查询和文档的相关性进行评分（带缓存）

    参数:
        query: 查询字符串
        doc: 文档内容

    返回:
        相关性得分 (0-10)
    """
    try:
        # 构建评分提示词
        prompt = f"""给定以下查询和文档片段，评估它们的相关性。
        评分标准：0分表示完全不相关，10分表示高度相关。
        只需返回一个0-10之间的整数分数，不要有任何其他解释。

        查询: {query}

        文档片段: {doc}

        相关性分数(0-10):"""

        # 调用本地LLM
        response = session.post(
            "http://localhost:11434/api/generate",
            json={
                "model": OLLAMA_MODEL_NAME,  # 通过环境变量 OLLAMA_MODEL_NAME 配置
                "prompt": prompt,
                "stream": False
            },
            timeout=180
        )

        # 提取得分
        result = response.json().get("response", "").strip()

        # 尝试解析为数字
        try:
            score = float(result)
            # 确保分数在0-10范围内
            score = max(0, min(10, score))
            return score
        except ValueError:
            # 如果无法解析为数字，尝试从文本中提取数字
            match = re.search(r'\b([0-9]|10)\b', result)
            if match:
                return float(match.group(1))
            else:
                # 默认返回中等相关性
                return 5.0

    except Exception as e:
        logging.error(f"LLM评分失败: {str(e)}")
        # 默认返回中等相关性
        return 5.0


def rerank_with_llm(query, docs, doc_ids, metadata_list, top_k=5):
    """
    使用LLM对检索结果进行重排序

    参数:
        query: 查询字符串
        docs: 文档内容列表
        doc_ids: 文档ID列表
        metadata_list: 元数据列表
        top_k: 返回结果数量

    返回:
        重排序后的结果列表
    """
    if not docs:
        return []

    results = []

    # 对每个文档进行评分
    for doc_id, doc, meta in zip(doc_ids, docs, metadata_list):
        # 获取LLM评分
        score = get_llm_relevance_score(query, doc)

        # 添加到结果列表
        results.append((doc_id, {
            'content': doc,
            'metadata': meta,
            'score': score / 10.0  # 归一化到0-1
        }))

    # 按得分排序
    results = sorted(results, key=lambda x: x[1]['score'], reverse=True)

    # 返回前K个结果
    return results[:top_k]

def rerank_results(query, docs, doc_ids, metadata_list, method=None, top_k=5):
    """
    对检索结果进行重排序

    参数:
        query: 查询字符串
        docs: 文档内容列表
        doc_ids: 文档ID列表
        metadata_list: 元数据列表
        method: 重排序方法 ("cross_encoder", "llm" 或 None)
        top_k: 返回结果数量

    返回:
        重排序后的结果
    """
    # 如果未指定方法，使用全局配置
    if method is None:
        method = RERANK_METHOD

    # 根据方法选择重排序函数
    if method == "llm":
        return rerank_with_llm(query, docs, doc_ids, metadata_list, top_k)
    elif method == "cross_encoder":
        return rerank_with_cross_encoder(query, docs, doc_ids, metadata_list, top_k)
    else:
        # 默认不进行重排序，按原始顺序返回
        return [(doc_id, {'content': doc, 'metadata': meta, 'score': 1.0 - idx / len(docs)})
                for idx, (doc_id, doc, meta) in enumerate(zip(doc_ids, docs, metadata_list))]

def stream_answer(question, enable_web_search=False, model_choice="siliconflow", progress=gr.Progress()):
    """改进的流式问答处理流程，支持联网搜索、混合检索和重排序，以及多种模型选择"""
    global faiss_index  # 确保可以访问
    try:
        # 检查向量数据库是否为空
        knowledge_base_exists = faiss_index is not None and faiss_index.ntotal > 0
        if not knowledge_base_exists:
            if not enable_web_search:
                yield "⚠️ 知识库为空，请先上传文档。", "遇到错误"
                return
            else:
                logging.warning("知识库为空，将仅使用网络搜索结果")

        progress(0.3, desc="执行递归检索...")
        # 使用递归检索获取更全面的答案上下文
        all_contexts, all_doc_ids, all_metadata = recursive_retrieval(
            initial_query=question,
            max_iterations=3,
            enable_web_search=enable_web_search,
            model_choice=model_choice
        )

        # 组合上下文，包含来源信息
        context_with_sources = []
        sources_for_conflict_detection = []

        # 使用检索到的结果构建上下文
        for doc, doc_id, metadata in zip(all_contexts, all_doc_ids, all_metadata):
            source_type = metadata.get('source', '本地文档')

            source_item = {
                'text': doc,
                'type': source_type
            }

            if source_type == 'web':
                url = metadata.get('url', '未知URL')
                title = metadata.get('title', '未知标题')
                context_with_sources.append(f"[网络来源: {title}] (URL: {url})\n{doc}")
                source_item['url'] = url
                source_item['title'] = title
            else:
                source = metadata.get('source', '未知来源')
                context_with_sources.append(f"[本地文档: {source}]\n{doc}")
                source_item['source'] = source

            sources_for_conflict_detection.append(source_item)

        # 检测矛盾
        conflict_detected = detect_conflicts(sources_for_conflict_detection)

        # 获取可信源
        if conflict_detected:
            credible_sources = [s for s in sources_for_conflict_detection
                                if s['type'] == 'web' and evaluate_source_credibility(s) > 0.7]

        context = "\n\n".join(context_with_sources)

        # 添加时间敏感检测
        time_sensitive = any(word in question for word in ["最新", "今年", "当前", "最近", "刚刚"])

        # 改进提示词模板，提高回答质量
        prompt_template = """作为一个专业的问答助手，你需要基于以下{context_type}回答用户问题。

提供的参考内容：
{context}

用户问题：{question}

请遵循以下回答原则：
1. 仅基于提供的参考内容回答问题，不要使用你自己的知识
2. 如果参考内容中没有足够信息，请坦诚告知你无法回答
3. 回答应该全面、准确、有条理，并使用适当的段落和结构
4. 请用中文回答
5. 在回答末尾标注信息来源{time_instruction}{conflict_instruction}

请现在开始回答："""

        prompt = prompt_template.format(
            context_type="本地文档和网络搜索结果" if enable_web_search and knowledge_base_exists else (
                "网络搜索结果" if enable_web_search else "本地文档"),
            context=context if context else (
                "网络搜索结果将用于回答。" if enable_web_search and not knowledge_base_exists else "知识库为空或未找到相关内容。"),
            question=question,
            time_instruction="，优先使用最新的信息" if time_sensitive and enable_web_search else "",
            conflict_instruction="，并明确指出不同来源的差异" if conflict_detected else ""
        )

        progress(0.7, desc="生成回答...")
        full_answer = ""

        # 根据模型选择使用不同的API
        if model_choice == "siliconflow":
            # 对于SiliconFlow API，不支持流式响应，所以一次性获取
            progress(0.8, desc="通过SiliconFlow API生成回答...")
            full_answer = call_siliconflow_api(prompt, temperature=0.7, max_tokens=1536)

            # 处理思维链
            if "<think>" in full_answer and "</think>" in full_answer:
                processed_answer = process_thinking_content(full_answer)
            else:
                processed_answer = full_answer

            yield processed_answer, "完成!"
        else:
            # 使用本地Ollama模型的流式响应
            response = session.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": OLLAMA_MODEL_NAME,
                    "prompt": prompt,
                    "stream": True
                },
                timeout=120,
                stream=True
            )

            for line in response.iter_lines():
                if line:
                    chunk = json.loads(line.decode()).get("response", "")
                    full_answer += chunk

                    # 检查是否有完整的思维链标签可以处理
                    if "<think>" in full_answer and "</think>" in full_answer:
                        # 需要确保完整收集一个思维链片段后再显示
                        processed_answer = process_thinking_content(full_answer)
                    else:
                        processed_answer = full_answer

                    yield processed_answer, "生成回答中..."

            # 处理最终输出，确保应用思维链处理
            final_answer = process_thinking_content(full_answer)
            yield final_answer, "完成!"

    except Exception as e:
        yield f"系统错误: {str(e)}", "遇到错误"


def query_answer(question, enable_web_search=False, model_choice="siliconflow", progress=gr.Progress()):
    """问答处理流程，支持联网搜索、混合检索和重排序，以及多种模型选择"""
    global faiss_index  # 确保可以访问
    try:
        logging.info(f"收到问题：{question}，联网状态：{enable_web_search}，模型选择：{model_choice}")

        # 检查向量数据库是否为空
        knowledge_base_exists = faiss_index is not None and faiss_index.ntotal > 0
        if not knowledge_base_exists:
            if not enable_web_search:
                return "⚠️ 知识库为空，请先上传文档。"
            else:
                logging.warning("知识库为空，将仅使用网络搜索结果")

        progress(0.3, desc="执行递归检索...")
        # 使用递归检索获取更全面的答案上下文
        all_contexts, all_doc_ids, all_metadata = recursive_retrieval(
            initial_query=question,
            max_iterations=3,
            enable_web_search=enable_web_search,
            model_choice=model_choice
        )

        # 组合上下文，包含来源信息
        context_with_sources = []
        sources_for_conflict_detection = []

        # 使用检索到的结果构建上下文
        for doc, doc_id, metadata in zip(all_contexts, all_doc_ids, all_metadata):
            source_type = metadata.get('source', '本地文档')

            source_item = {
                'text': doc,
                'type': source_type
            }

            if source_type == 'web':
                url = metadata.get('url', '未知URL')
                title = metadata.get('title', '未知标题')
                context_with_sources.append(f"[网络来源: {title}] (URL: {url})\n{doc}")
                source_item['url'] = url
                source_item['title'] = title
            else:
                source = metadata.get('source', '未知来源')
                context_with_sources.append(f"[本地文档: {source}]\n{doc}")
                source_item['source'] = source

            sources_for_conflict_detection.append(source_item)

        # 检测矛盾
        conflict_detected = detect_conflicts(sources_for_conflict_detection)

        # 获取可信源
        if conflict_detected:
            credible_sources = [s for s in sources_for_conflict_detection
                                if s['type'] == 'web' and evaluate_source_credibility(s) > 0.7]

        context = "\n\n".join(context_with_sources)

        # 添加时间敏感检测
        time_sensitive = any(word in question for word in ["最新", "今年", "当前", "最近", "刚刚"])

        # 改进提示词模板，提高回答质量
        prompt_template = """作为一个专业的问答助手，你需要基于以下{context_type}回答用户问题。

提供的参考内容：
{context}

用户问题：{question}

请遵循以下回答原则：
1. 仅基于提供的参考内容回答问题，不要使用你自己的知识
2. 如果参考内容中没有足够信息，请坦诚告知你无法回答
3. 回答应该全面、准确、有条理，并使用适当的段落和结构
4. 请用中文回答
5. 在回答末尾标注信息来源{time_instruction}{conflict_instruction}

请现在开始回答："""

        prompt = prompt_template.format(
            context_type="本地文档和网络搜索结果" if enable_web_search and knowledge_base_exists else (
                "网络搜索结果" if enable_web_search else "本地文档"),
            context=context if context else (
                "网络搜索结果将用于回答。" if enable_web_search and not knowledge_base_exists else "知识库为空或未找到相关内容。"),
            question=question,
            time_instruction="，优先使用最新的信息" if time_sensitive and enable_web_search else "",
            conflict_instruction="，并明确指出不同来源的差异" if conflict_detected else ""
        )

        progress(0.8, desc="生成回答...")

        # 根据模型选择使用不同的API
        if model_choice == "siliconflow":
            # 使用SiliconFlow API
            result = call_siliconflow_api(prompt, temperature=0.7, max_tokens=1536)

            # 处理思维链
            processed_result = process_thinking_content(result)
            return processed_result
        else:
            # 使用本地Ollama
            response = session.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": OLLAMA_MODEL_NAME,
                    "prompt": prompt,
                    "stream": False
                },
                timeout=180,  # 延长到3分钟
                headers={'Connection': 'close'}  # 添加连接头
            )
            response.raise_for_status()  # 检查HTTP状态码

            progress(1.0, desc="完成!")
            # 确保返回字符串并处理空值
            result = response.json()
            return process_thinking_content(str(result.get("response", "未获取到有效回答")))

    except json.JSONDecodeError:
        return "响应解析失败，请重试"
    except KeyError:
        return "响应格式异常，请检查模型服务"
    except Exception as e:
        progress(1.0, desc="遇到错误")  # 确保进度条完成
        return f"系统错误: {str(e)}"


def process_thinking_content(text):
    """处理包含<think>标签的内容，将其转换为Markdown格式"""
    # 检查输入是否为有效文本
    if text is None:
        return ""

    # 确保输入是字符串
    if not isinstance(text, str):
        try:
            processed_text = str(text)
        except:
            return "无法处理的内容格式"
    else:
        processed_text = text

    # 处理思维链标签
    try:
        while "<think>" in processed_text and "</think>" in processed_text:
            start_idx = processed_text.find("<think>")
            end_idx = processed_text.find("</think>")
            if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
                thinking_content = processed_text[start_idx + 7:end_idx]
                before_think = processed_text[:start_idx]
                after_think = processed_text[end_idx + 8:]

                # 使用可折叠详情框显示思维链
                processed_text = before_think + "\n\n<details>\n<summary>思考过程（点击展开）</summary>\n\n" + thinking_content + "\n\n</details>\n\n" + after_think

        # 处理其他HTML标签，但保留details和summary标签
        processed_html = []
        i = 0
        while i < len(processed_text):
            if processed_text[i:i + 8] == "<details" or processed_text[i:i + 9] == "</details" or \
                    processed_text[i:i + 8] == "<summary" or processed_text[i:i + 9] == "</summary":
                # 保留这些标签
                tag_end = processed_text.find(">", i)
                if tag_end != -1:
                    processed_html.append(processed_text[i:tag_end + 1])
                    i = tag_end + 1
                    continue

            if processed_text[i] == "<":
                processed_html.append("&lt;")
            elif processed_text[i] == ">":
                processed_html.append("&gt;")
            else:
                processed_html.append(processed_text[i])
            i += 1

        processed_text = "".join(processed_html)
    except Exception as e:
        logging.error(f"处理思维链内容时出错: {str(e)}")
        # 出错时至少返回原始文本，但确保安全处理HTML标签
        try:
            return text.replace("<", "&lt;").replace(">", "&gt;")
        except:
            return "处理内容时出错"

    return processed_text


def call_siliconflow_api(prompt, temperature=0.7, max_tokens=1024):
    """
    调用SiliconFlow API获取回答

    Args:
        prompt: 提示词
        temperature: 温度参数
        max_tokens: 最大生成token数

    Returns:
        生成的回答文本和思维链内容
    """
    # 检查是否配置了SiliconFlow API密钥
    if not SILICONFLOW_API_KEY:
        logging.error("未设置 SILICONFLOW_API_KEY 环境变量。请在.env文件中设置您的 API 密钥。")
        return "错误：未配置 SiliconFlow API 密钥。", ""

    try:
        payload = {
            "model": SILICONFLOW_MODEL_NAME,  # 通过环境变量 SILICONFLOW_MODEL_NAME 配置
            "messages": [
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            "stream": False,
            "max_tokens": max_tokens,
            "stop": None,
            "temperature": temperature,
            "top_p": 0.7,
            "top_k": 50,
            "frequency_penalty": 0.5,
            "n": 1,
            "response_format": {"type": "text"}
        }

        headers = {
            "Authorization": f"Bearer {SILICONFLOW_API_KEY.strip()}",  # 从环境变量获取密钥并去除空格
            "Content-Type": "application/json; charset=utf-8"  # 明确指定编码
        }

        # 手动将payload编码为UTF-8 JSON字符串
        json_payload = json.dumps(payload, ensure_ascii=False).encode('utf-8')

        response = requests.post(
            SILICONFLOW_API_URL,
            data=json_payload,  # 通过data参数发送编码后的JSON
            headers=headers,
            timeout=180  # 延长超时时间到3分钟
        )

        response.raise_for_status()
        result = response.json()

        # 提取回答内容和思维链
        if "choices" in result and len(result["choices"]) > 0:
            message = result["choices"][0]["message"]
            content = message.get("content", "")
            reasoning = message.get("reasoning_content", "")

            # 如果有思维链，则添加特殊标记，以便前端处理
            if reasoning:
                # 添加思维链标记
                full_response = f"{content}<think>{reasoning}</think>"
                return full_response
            else:
                return content
        else:
            return "API返回结果格式异常，请检查"

    except requests.exceptions.RequestException as e:
        logging.error(f"调用SiliconFlow API时出错: {str(e)}")
        return f"调用API时出错: {str(e)}"
    except json.JSONDecodeError:
        logging.error("SiliconFlow API返回非JSON响应")
        return "API响应解析失败"
    except Exception as e:
        logging.error(f"调用SiliconFlow API时发生未知错误: {str(e)}")
        return f"发生未知错误: {str(e)}"

# 合并语义搜索和BM25搜索结果
def hybrid_merge(semantic_results, bm25_results, alpha=0.7):
    """
    合并语义搜索和BM25搜索结果

    参数:
        semantic_results: 向量检索结果 (字典格式，包含ids, documents, metadatas)
        bm25_results: BM25检索结果 (字典列表，包含id, score, content)
        alpha: 语义搜索权重 (0-1)

    返回:
        合并后的结果列表 [(doc_id, {'score': score, 'content': content, 'metadata': metadata}), ...]
    """
    merged_dict = {}
    global faiss_metadatas_map  # Ensure we can access the global map

    # 处理语义搜索结果
    if (semantic_results and
            isinstance(semantic_results.get('documents'), list) and len(semantic_results['documents']) > 0 and
            isinstance(semantic_results.get('metadatas'), list) and len(semantic_results['metadatas']) > 0 and
            isinstance(semantic_results.get('ids'), list) and len(semantic_results['ids']) > 0 and
            isinstance(semantic_results['documents'][0], list) and
            isinstance(semantic_results['metadatas'][0], list) and
            isinstance(semantic_results['ids'][0], list) and
            len(semantic_results['documents'][0]) == len(semantic_results['metadatas'][0]) == len(
                semantic_results['ids'][0])):

        num_results = len(semantic_results['documents'][0])
        # Assuming semantic_results are already ordered by relevance (higher is better)
        # A simple rank-based score, can be replaced if actual scores/distances are available and preferred
        for i, (doc_id, doc, meta) in enumerate(
                zip(semantic_results['ids'][0], semantic_results['documents'][0], semantic_results['metadatas'][0])):
            score = 1.0 - (i / max(1, num_results))  # Higher rank (smaller i) gets higher score
            merged_dict[doc_id] = {
                'score': alpha * score,
                'content': doc,
                'metadata': meta
            }
    else:
        logging.warning(
            "Semantic results are missing, have an unexpected format, or are empty. Skipping semantic part in hybrid merge.")

    # 处理BM25结果
    if not bm25_results:
        return sorted(merged_dict.items(), key=lambda x: x[1]['score'], reverse=True)

    valid_bm25_scores = [r['score'] for r in bm25_results if isinstance(r, dict) and 'score' in r]
    max_bm25_score = max(valid_bm25_scores) if valid_bm25_scores else 1.0

    for result in bm25_results:
        if not (isinstance(result, dict) and 'id' in result and 'score' in result and 'content' in result):
            logging.warning(f"Skipping invalid BM25 result item: {result}")
            continue

        doc_id = result['id']
        # Normalize BM25 score
        normalized_score = result['score'] / max_bm25_score if max_bm25_score > 0 else 0

        if doc_id in merged_dict:
            merged_dict[doc_id]['score'] += (1 - alpha) * normalized_score
        else:
            metadata = faiss_metadatas_map.get(doc_id, {})  # Get metadata from our global map
            merged_dict[doc_id] = {
                'score': (1 - alpha) * normalized_score,
                'content': result['content'],
                'metadata': metadata
            }

    merged_results = sorted(merged_dict.items(), key=lambda x: x[1]['score'], reverse=True)
    return merged_results

def update_bm25_index():
    """更新BM25索引，从内存中的映射加载所有文档"""
    global faiss_contents_map, faiss_id_order_for_index
    try:
        # Use the ordered list of IDs to ensure consistency
        doc_ids = faiss_id_order_for_index
        if not doc_ids:
            logging.warning("没有可索引的文档 (FAISS ID列表为空)")
            BM25_MANAGER.clear()
            return False

        # Retrieve documents in the correct order
        documents = [faiss_contents_map.get(doc_id, "") for doc_id in doc_ids]

        # Filter out any potential empty documents if necessary, though map access should be safe
        valid_docs_with_ids = [(doc_id, doc) for doc_id, doc in zip(doc_ids, documents) if doc]
        if not valid_docs_with_ids:
            logging.warning("没有有效的文档内容可用于BM25索引")
            BM25_MANAGER.clear()
            return False

        # Separate IDs and documents again for building the index
        final_doc_ids = [item[0] for item in valid_docs_with_ids]
        final_documents = [item[1] for item in valid_docs_with_ids]

        BM25_MANAGER.build_index(final_documents, final_doc_ids)
        logging.info(f"BM25索引更新完成，共索引 {len(final_doc_ids)} 个文档")
        return True
    except Exception as e:
        logging.error(f"更新BM25索引失败: {str(e)}")
        return False

def get_system_models_info():
    """返回系统使用的各种模型信息"""
    models_info = {
        "嵌入模型": "all-MiniLM-L6-v2",
        "分块方法": "RecursiveCharacterTextSplitter (chunk_size=800, overlap=150)",
        "检索方法": "向量检索 + BM25混合检索 (α=0.7)",
        "重排序模型": "交叉编码器 (sentence-transformers/distiluse-base-multilingual-cased-v2)",
        "生成模型(Ollama)": OLLAMA_MODEL_NAME,
        "生成模型(SiliconFlow)": SILICONFLOW_MODEL_NAME,
        "分词工具": "jieba (中文分词)"
    }
    return models_info


# 修改全局缓存变量为字典格式，便于通过ID快速查找
chunk_data_cache = {}  # 格式: {chunk_id: chunk_data}
def get_document_chunks(progress=gr.Progress()):
    """获取文档分块结果用于可视化"""
    global faiss_contents_map, faiss_metadatas_map, faiss_id_order_for_index
    global chunk_data_cache  # 声明使用全局缓存

    try:
        progress(0.1, desc="正在从内存加载数据...")

        # 清空旧缓存
        chunk_data_cache.clear()

        if not faiss_id_order_for_index:
            return [], "知识库中没有文档，请先上传并处理文档。"

        progress(0.3, desc="正在组织分块数据...")

        # 按原始处理顺序组织数据
        table_data = []
        for idx, chunk_id in enumerate(faiss_id_order_for_index):
            content = faiss_contents_map.get(chunk_id, "")
            meta = faiss_metadatas_map.get(chunk_id, {})

            if not content:
                continue

            # 构建分块数据对象
            chunk_data = {
                "row_id": idx,  # 表格行号
                "chunk_id": chunk_id,
                "source": meta.get("source", "未知来源"),
                "content": content,
                "preview": content[:200] + "..." if len(content) > 200 else content,
                "char_count": len(content),
                "token_count": len(list(jieba.cut(content)))
            }

            # 添加到缓存和表格数据
            chunk_data_cache[idx] = chunk_data  # 用行号作为键
            table_data.append([
                chunk_data["source"],
                f"{idx + 1}/{len(faiss_id_order_for_index)}",
                chunk_data["char_count"],
                chunk_data["token_count"],
                chunk_data["preview"]
            ])

        progress(1.0, desc="数据加载完成!")
        return table_data, f"共 {len(table_data)} 个文本块"

    except Exception as e:
        chunk_data_cache.clear()
        return [], f"获取分块数据失败: {str(e)}"


def show_chunk_details(evt: gr.SelectData):
    """显示选中分块的详细内容"""
    try:
        if not evt.index or evt.index[0] is None:
            return "未选择有效行"

        row_idx = evt.index[0]  # 获取行索引

        # 从缓存获取数据
        selected_chunk = chunk_data_cache.get(row_idx)
        if not selected_chunk:
            return "未找到对应的分块数据"

        # 格式化显示详情
        detail = f"""
        [来源] {selected_chunk['source']}
        [ID] {selected_chunk['chunk_id']}
        [字符数] {selected_chunk['char_count']}
        [分词数] {selected_chunk['token_count']}
        ----------------------------
        {selected_chunk['content']}
        """
        return detail

    except Exception as e:
        return f"加载分块详情失败: {str(e)}"

# 修改布局部分，添加一个新的标签页
# 修改布局部分，添加真实系统监控功能
with gr.Blocks(
        title="本地RAG问答系统",
        css="""
    /* 全局主题变量 */
    :root[data-theme="light"] {
        --text-color: #2c3e50;
        --bg-color: #ffffff;
        --panel-bg: #f8f9fa;
        --border-color: #e9ecef;
        --success-color: #4CAF50;
        --error-color: #f44336;
        --primary-color: #2196F3;
        --secondary-bg: #ffffff;
        --hover-color: #e9ecef;
        --chat-user-bg: #e3f2fd;
        --chat-assistant-bg: #f5f5f5;
        --tech-blue: #0d47a1;
        --tech-purple: #7b1fa2;
        --tech-cyan: #00bcd4;
    }

    :root[data-theme="dark"] {
        --text-color: #e0e0e0;
        --bg-color: #1a1a1a;
        --panel-bg: #2d2d2d;
        --border-color: #404040;
        --success-color: #81c784;
        --error-color: #e57373;
        --primary-color: #64b5f6;
        --secondary-bg: #2d2d2d;
        --hover-color: #404040;
        --chat-user-bg: #1e3a5f;
        --chat-assistant-bg: #2d2d2d;
        --tech-blue: #1e88e5;
        --tech-purple: #9c27b0;
        --tech-cyan: #00e5ff;
    }

    /* 全局样式 */
    body {
        font-family: 'Roboto', 'Segoe UI', sans-serif;
        margin: 0;
        padding: 0;
        overflow-x: hidden;
        width: 100vw;
        height: 100vh;
        background: linear-gradient(135deg, var(--bg-color) 0%, #1a1a2e 100%);
    }

    .gradio-container {
        max-width: 100% !important;
        width: 100% !important;
        margin: 0 !important;
        padding: 0 1% !important;
        color: var(--text-color);
        background-color: transparent;
        min-height: 100vh;
    }

    /* 确保标签内容撑满 */
    .tabs.svelte-710i53 {
        margin: 0 !important;
        padding: 0 !important;
        width: 100% !important;
    }

    /* 主题切换按钮 */
    .theme-toggle {
        position: fixed;
        top: 20px;
        right: 20px;
        z-index: 1000;
        padding: 8px 16px;
        border-radius: 20px;
        border: 1px solid var(--border-color);
        background: var(--panel-bg);
        color: var(--text-color);
        cursor: pointer;
        transition: all 0.3s ease;
        font-size: 14px;
        display: flex;
        align-items: center;
        gap: 8px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    }

    .theme-toggle:hover {
        background: var(--hover-color);
        transform: translateY(-2px);
        box-shadow: 0 6px 8px rgba(0,0,0,0.15);
    }

    /* 面板样式 */
    .left-panel {
        padding-right: 20px;
        border-right: 1px solid var(--border-color);
        background: rgba(30, 30, 46, 0.7);
        width: 100%;
        border-radius: 12px;
        backdrop-filter: blur(10px);
        box-shadow: 0 8px 32px rgba(0,0,0,0.2);
    }

    .right-panel {
        height: 100vh;
        background: rgba(30, 30, 46, 0.7);
        width: 100%;
        border-radius: 12px;
        backdrop-filter: blur(10px);
        box-shadow: 0 8px 32px rgba(0,0,0,0.2);
    }

    /* 文件列表样式 */
    .file-list {
        margin-top: 10px;
        padding: 12px;
        background: rgba(45, 45, 70, 0.6);
        border-radius: 8px;
        font-size: 14px;
        line-height: 1.6;
        border: 1px solid rgba(100, 100, 150, 0.3);
    }

    /* 答案框样式 */
    .answer-box {
        min-height: 500px !important;
        background: rgba(45, 45, 70, 0.6);
        border-radius: 8px;
        padding: 16px;
        font-size: 15px;
        line-height: 1.6;
        border: 1px solid rgba(100, 100, 150, 0.3);
    }

    /* 输入框样式 */
    textarea {
        background: rgba(45, 45, 70, 0.6) !important;
        color: var(--text-color) !important;
        border: 1px solid rgba(100, 100, 150, 0.3) !important;
        border-radius: 8px !important;
        padding: 12px !important;
        font-size: 14px !important;
        transition: all 0.3s ease;
    }

    textarea:focus {
        border-color: var(--tech-cyan) !important;
        box-shadow: 0 0 0 2px rgba(0, 188, 212, 0.2);
    }

    /* 按钮样式 */
    button.primary {
        background: linear-gradient(135deg, var(--tech-blue) 0%, var(--tech-purple) 100%) !important;
        color: white !important;
        border-radius: 8px !important;
        padding: 10px 20px !important;
        font-weight: 500 !important;
        transition: all 0.3s ease !important;
        border: none;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    }

    button.primary:hover {
        opacity: 0.9;
        transform: translateY(-2px);
        box-shadow: 0 6px 8px rgba(0,0,0,0.15);
    }

    /* 标题和文本样式 */
    h1, h2, h3 {
        color: var(--text-color) !important;
        font-weight: 600 !important;
    }

    .footer-note {
        color: var(--text-color);
        opacity: 0.8;
        font-size: 13px;
        margin-top: 12px;
    }

    /* 加载和进度样式 */
    #loading, .progress-text {
        color: var(--text-color);
    }

    /* 聊天记录样式 */
    .chat-container {
        border: 1px solid rgba(100, 100, 150, 0.3);
        border-radius: 8px;
        margin-bottom: 16px;
        max-height: 80vh;
        height: 80vh !important;
        overflow-y: auto;
        background: rgba(45, 45, 70, 0.6);
    }

    .chat-message {
        padding: 12px 16px;
        margin: 8px;
        border-radius: 8px;
        font-size: 14px;
        line-height: 1.5;
        position: relative;
        overflow: hidden;
    }

    .chat-message.user {
        background: linear-gradient(135deg, rgba(30, 58, 95, 0.8) 0%, rgba(30, 30, 70, 0.8) 100%);
        margin-left: 32px;
        border-top-right-radius: 4px;
        border-left: 3px solid var(--tech-cyan);
    }

    .chat-message.assistant {
        background: linear-gradient(135deg, rgba(45, 45, 70, 0.8) 0%, rgba(30, 30, 50, 0.8) 100%);
        margin-right: 32px;
        border-top-left-radius: 4px;
        border-right: 3px solid var(--tech-purple);
    }

    .chat-message .timestamp {
        font-size: 12px;
        color: var(--text-color);
        opacity: 0.7;
        margin-bottom: 4px;
    }

    .chat-message .content {
        white-space: pre-wrap;
    }

    /* 按钮组样式 */
    .button-row {
        display: flex;
        gap: 8px;
        margin-top: 8px;
    }

    .clear-button {
        background: linear-gradient(135deg, #f44336 0%, #c62828 100%) !important;
    }

    /* API配置提示样式 */
    .api-info {
        margin-top: 10px;
        padding: 10px;
        border-radius: 5px;
        background: rgba(45, 45, 70, 0.6);
        border: 1px solid rgba(100, 100, 150, 0.3);
    }

    /* 新增: 数据可视化卡片样式 */
    .model-card {
        background: rgba(45, 45, 70, 0.6);
        border-radius: 8px;
        padding: 16px;
        border: 1px solid rgba(100, 100, 150, 0.3);
        margin-bottom: 16px;
    }

    .model-card h3 {
        margin-top: 0;
        border-bottom: 1px solid rgba(100, 100, 150, 0.3);
        padding-bottom: 8px;
        color: var(--tech-cyan);
    }

    .model-item {
        display: flex;
        margin-bottom: 8px;
    }

    .model-item .label {
        flex: 1;
        font-weight: 500;
        color: var(--tech-cyan);
    }

    .model-item .value {
        flex: 2;
    }

    /* 数据表格样式 */
    .chunk-table {
        border-radius: 8px;
        overflow: hidden;
        border: 1px solid rgba(100, 100, 150, 0.3);
        background: rgba(45, 45, 70, 0.6);
    }

    .chunk-table th, .chunk-table td {
        border: 1px solid rgba(100, 100, 150, 0.3);
        padding: 8px;
    }

    .chunk-detail-box {
        min-height: 200px;
        padding: 16px;
        background: rgba(45, 45, 70, 0.6);
        border-radius: 8px;
        border: 1px solid rgba(100, 100, 150, 0.3);
        font-family: monospace;
        white-space: pre-wrap;
        overflow-y: auto;
    }

    /* 新增: 系统监控面板样式 */
    .monitor-panel {
        background: rgba(30, 30, 46, 0.7);
        border-radius: 12px;
        padding: 20px;
        margin-bottom: 20px;
        backdrop-filter: blur(10px);
        box-shadow: 0 8px 32px rgba(0,0,0,0.2);
        border: 1px solid rgba(100, 100, 150, 0.3);
    }

    .monitor-header {
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-bottom: 15px;
        padding-bottom: 10px;
        border-bottom: 1px solid rgba(100, 100, 150, 0.3);
    }

    .monitor-title {
        font-size: 18px;
        font-weight: 600;
        color: var(--tech-cyan);
    }

    .monitor-refresh {
        background: transparent;
        border: none;
        color: var(--tech-cyan);
        cursor: pointer;
        font-size: 14px;
        display: flex;
        align-items: center;
        gap: 5px;
    }

    .monitor-grid {
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
        gap: 20px;
    }

    .metric-card {
        background: rgba(45, 45, 70, 0.6);
        border-radius: 10px;
        padding: 15px;
        border: 1px solid rgba(100, 100, 150, 0.3);
    }

    .metric-title {
        font-size: 14px;
        margin-bottom: 10px;
        color: var(--tech-cyan);
    }

    .metric-value {
        font-size: 24px;
        font-weight: 700;
        margin-bottom: 5px;
    }

    .metric-trend {
        font-size: 12px;
        color: #4CAF50;
    }

    .metric-trend.negative {
        color: #f44336;
    }

    .metric-chart {
        height: 100px;
        margin-top: 10px;
        position: relative;
    }

    .chart-bar {
        position: absolute;
        bottom: 0;
        width: 8px;
        background: var(--tech-cyan);
        border-radius: 4px 4px 0 0;
        transition: height 0.5s ease;
    }

    .log-container {
        max-height: 300px;
        overflow-y: auto;
        background: rgba(20, 20, 35, 0.8);
        border-radius: 8px;
        padding: 15px;
        font-family: monospace;
        font-size: 13px;
        line-height: 1.5;
    }

    .log-entry {
        margin-bottom: 8px;
        padding-bottom: 8px;
        border-bottom: 1px solid rgba(100, 100, 150, 0.2);
    }

    .log-time {
        color: var(--tech-cyan);
        margin-right: 10px;
    }

    .log-info {
        color: #4CAF50;
    }

    .log-warning {
        color: #FFC107;
    }

    .log-error {
        color: #f44336;
    }

    /* 新增: 科技感装饰元素 */
    .tech-grid {
        position: fixed;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        pointer-events: none;
        z-index: -1;
        opacity: 0.05;
    }

    .grid-line {
        position: absolute;
        background: var(--tech-cyan);
    }

    .grid-horizontal {
        width: 100%;
        height: 1px;
        top: 0;
        left: 0;
    }

    .grid-vertical {
        height: 100%;
        width: 1px;
        top: 0;
        left: 0;
    }

    /* 新增: 霓虹效果 */
    .neon-text {
        text-shadow: 0 0 5px var(--tech-cyan), 0 0 10px var(--tech-cyan), 0 0 15px var(--tech-purple);
    }

    /* 新增: 进度条样式 */
    .progress-container {
        width: 100%;
        background: rgba(255,255,255,0.1);
        border-radius: 10px;
        margin: 10px 0;
    }

    .progress-bar {
        height: 8px;
        border-radius: 10px;
        background: linear-gradient(90deg, var(--tech-cyan), var(--tech-purple));
        transition: width 0.3s ease;
    }
    """
) as demo:
    # 添加科技感网格背景
    gr.HTML("""
    <div class="tech-grid">
        <div class="grid-line grid-horizontal" style="top: 0;"></div>
        <div class="grid-line grid-horizontal" style="top: 20%;"></div>
        <div class="grid-line grid-horizontal" style="top: 40%;"></div>
        <div class="grid-line grid-horizontal" style="top: 60%;"></div>
        <div class="grid-line grid-horizontal" style="top: 80%;"></div>
        <div class="grid-line grid-horizontal" style="top: 100%;"></div>

        <div class="grid-line grid-vertical" style="left: 0;"></div>
        <div class="grid-line grid-vertical" style="left: 20%;"></div>
        <div class="grid-line grid-vertical" style="left: 40%;"></div>
        <div class="grid-line grid-vertical" style="left: 60%;"></div>
        <div class="grid-line grid-vertical" style="left: 80%;"></div>
        <div class="grid-line grid-vertical" style="left: 100%;"></div>
    </div>
    """)

    gr.Markdown("# 🧠 <span class='neon-text'>智能文档问答系统</span>")

    with gr.Tabs() as tabs:
        # 第一个选项卡：问答对话
        with gr.TabItem("💬 问答对话"):
            with gr.Row(equal_height=True):
                # 左侧操作面板 - 调整比例为合适的大小
                with gr.Column(scale=5, elem_classes="left-panel"):
                    gr.Markdown("## 📂 文档处理区")
                    with gr.Group():
                        # 在 Gradio UI 部分修改文件上传组件
                        file_input = gr.File(
                            label="上传文档 (支持PDF, Word, Excel, PPT, TXT, Markdown等)",
                            file_types=[".pdf", ".txt", ".docx", ".xlsx", ".xls", ".pptx", ".md"],
                            file_count="multiple"
                        )
                        upload_btn = gr.Button("🚀 开始处理", variant="primary")
                        upload_status = gr.Textbox(
                            label="处理状态",
                            interactive=False,
                            lines=2
                        )
                        file_list = gr.Textbox(
                            label="已处理文件",
                            interactive=False,
                            lines=3,
                            elem_classes="file-list"
                        )

                    # 将问题输入区移至左侧面板底部
                    gr.Markdown("## ❓ 输入问题")
                    with gr.Group():
                        question_input = gr.Textbox(
                            label="输入问题",
                            lines=3,
                            placeholder="请输入您的问题...",
                            elem_id="question-input"
                        )
                        with gr.Row():
                            # 添加联网开关
                            web_search_checkbox = gr.Checkbox(
                                label="启用联网搜索",
                                value=False,
                                info="打开后将同时搜索网络内容（需配置SERPAPI_KEY）"
                            )

                            # 添加模型选择下拉框
                            model_choice = gr.Dropdown(
                                choices=["ollama", "siliconflow"],
                                value="ollama",
                                label="模型选择",
                                info="选择使用本地模型或云端模型"
                            )

                        with gr.Row():
                            ask_btn = gr.Button("🔍 开始提问", variant="primary", scale=2)
                            clear_btn = gr.Button("🗑️ 清空对话", variant="secondary", elem_classes="clear-button",
                                                  scale=1)

                    # 添加API配置提示信息
                    api_info = gr.HTML(
                        """
                        <div class="api-info" style="margin-top:10px;padding:10px;border-radius:5px;background:var(--panel-bg);border:1px solid var(--border-color);">
                            <p>📢 <strong>功能说明：</strong></p>
                            <p>1. <strong>联网搜索</strong>：%s</p>
                            <p>2. <strong>模型选择</strong>：当前使用 <strong>%s</strong> %s</p>
                        </div>
                        """
                    )

                # 右侧对话区 - 调整比例
                with gr.Column(scale=7, elem_classes="right-panel"):
                    gr.Markdown("## 📝 对话记录")

                    # 对话记录显示区
                    chatbot = gr.Chatbot(
                        label="对话历史",
                        height=600,  # 增加高度
                        elem_classes="chat-container",
                        show_label=False
                    )

                    status_display = gr.HTML("", elem_id="status-display")
                    gr.Markdown("""
                    <div class="footer-note">
                        *回答生成可能需要1-2分钟，请耐心等待<br>
                        *支持多轮对话，可基于前文继续提问
                    </div>
                    """)

        # 第二个选项卡：分块可视化
        with gr.TabItem("📊 分块可视化"):
            with gr.Row():
                with gr.Column(scale=1):
                    gr.Markdown("## 💡 系统模型信息")

                    # 显示系统模型信息卡片
                    models_info = get_system_models_info()
                    with gr.Group(elem_classes="model-card"):
                        gr.Markdown("### 核心模型与技术")

                        for key, value in models_info.items():
                            with gr.Row():
                                gr.Markdown(f"**{key}**:", elem_classes="label")
                                gr.Markdown(f"{value}", elem_classes="value")

                with gr.Column(scale=2):
                    gr.Markdown("## 📄 文档分块统计")
                    refresh_chunks_btn = gr.Button("🔄 刷新分块数据", variant="primary")
                    chunks_status = gr.Markdown("点击按钮查看分块统计")

            # 分块数据表格和详情
            with gr.Row():
                chunks_data = gr.Dataframe(
                    headers=["来源", "序号", "字符数", "分词数", "内容预览"],
                    elem_classes="chunk-table",
                    interactive=False,
                    wrap=True,
                    row_count=(10, "dynamic")
                )

            with gr.Row():
                chunk_detail_text = gr.Textbox(
                    label="分块详情",
                    placeholder="点击表格中的行查看完整内容...",
                    lines=8,
                    elem_classes="chunk-detail-box"
                )

            gr.Markdown("""
            <div class="footer-note">
                * 点击表格中的行可查看该分块的完整内容<br>
                * 分词数表示使用jieba分词后的token数量
            </div>
            """)

        # 新增第三个选项卡：系统监控
        with gr.TabItem("📈 系统监控"):
            with gr.Column():
                # 系统资源监控面板
                with gr.Group(elem_classes="monitor-panel"):
                    with gr.Row():
                        gr.Markdown("## 🖥️ 系统资源监控", elem_classes="monitor-title")
                        refresh_monitor_btn = gr.Button("🔄 刷新数据", variant="primary", elem_classes="monitor-refresh")

                    with gr.Row(elem_classes="monitor-grid"):
                        # CPU使用率
                        with gr.Column():
                            cpu_card = gr.Group(elem_classes="metric-card")
                            with cpu_card:
                                gr.Markdown("CPU使用率", elem_classes="metric-title")
                                cpu_value = gr.Markdown("加载中...", elem_classes="metric-value")
                                cpu_progress = gr.HTML("""
                                    <div class="progress-container">
                                        <div class="progress-bar" style="width: 0%"></div>
                                    </div>
                                """)
                                cpu_info = gr.Markdown("核心数: 加载中...", elem_classes="metric-trend")

                        # 内存使用
                        with gr.Column():
                            memory_card = gr.Group(elem_classes="metric-card")
                            with memory_card:
                                gr.Markdown("内存使用", elem_classes="metric-title")
                                memory_value = gr.Markdown("加载中...", elem_classes="metric-value")
                                memory_progress = gr.HTML("""
                                    <div class="progress-container">
                                        <div class="progress-bar" style="width: 0%"></div>
                                    </div>
                                """)
                                memory_info = gr.Markdown("总内存: 加载中...", elem_classes="metric-trend")

                        # 磁盘空间
                        with gr.Column():
                            disk_card = gr.Group(elem_classes="metric-card")
                            with disk_card:
                                gr.Markdown("磁盘空间", elem_classes="metric-title")
                                disk_value = gr.Markdown("加载中...", elem_classes="metric-value")
                                disk_progress = gr.HTML("""
                                    <div class="progress-container">
                                        <div class="progress-bar" style="width: 0%"></div>
                                    </div>
                                """)
                                disk_info = gr.Markdown("总空间: 加载中...", elem_classes="metric-trend")

                        # 网络流量
                        with gr.Column():
                            network_card = gr.Group(elem_classes="metric-card")
                            with network_card:
                                gr.Markdown("网络流量", elem_classes="metric-title")
                                network_value = gr.Markdown("加载中...", elem_classes="metric-value")
                                network_info = gr.Markdown("上传/下载: 0 KB/s", elem_classes="metric-trend")

                # 性能指标面板
                with gr.Group(elem_classes="monitor-panel"):
                    gr.Markdown("## ⚡ 性能指标", elem_classes="monitor-title")

                    with gr.Row(elem_classes="monitor-grid"):
                        # 响应时间
                        with gr.Column():
                            latency_card = gr.Group(elem_classes="metric-card")
                            with latency_card:
                                gr.Markdown("平均响应时间", elem_classes="metric-title")
                                latency_value = gr.Markdown("0 ms", elem_classes="metric-value")
                                latency_info = gr.Markdown("历史记录: 加载中...", elem_classes="metric-trend")

                        # 请求速率
                        with gr.Column():
                            request_card = gr.Group(elem_classes="metric-card")
                            with request_card:
                                gr.Markdown("请求统计", elem_classes="metric-title")
                                request_value = gr.Markdown("总请求: 0", elem_classes="metric-value")
                                request_info = gr.Markdown("成功/失败: 0/0", elem_classes="metric-trend")

                        # 向量数据库
                        with gr.Column():
                            vector_db_card = gr.Group(elem_classes="metric-card")
                            with vector_db_card:
                                gr.Markdown("向量数据库", elem_classes="metric-title")
                                vector_db_value = gr.Markdown("分块数: 0", elem_classes="metric-value")
                                vector_db_info = gr.Markdown("向量数: 0", elem_classes="metric-trend")

                        # 模型状态
                        with gr.Column():
                            model_status_card = gr.Group(elem_classes="metric-card")
                            with model_status_card:
                                gr.Markdown("模型状态", elem_classes="metric-title")
                                model_status_value = gr.Markdown("状态: 未知", elem_classes="metric-value")
                                model_status_info = gr.Markdown("连接: 检查中...", elem_classes="metric-trend")

                # 系统日志面板
                with gr.Group(elem_classes="monitor-panel"):
                    gr.Markdown("## 📝 系统日志", elem_classes="monitor-title")

                    # 日志筛选选项
                    with gr.Row():
                        log_level = gr.Dropdown(
                            choices=["所有级别", "信息", "警告", "错误"],
                            value="所有级别",
                            label="日志级别"
                        )
                        log_search = gr.Textbox(
                            label="搜索日志",
                            placeholder="输入关键词搜索..."
                        )
                        clear_logs_btn = gr.Button("🗑️ 清空日志", variant="secondary")

                    # 日志显示区域
                    log_display = gr.HTML("", elem_classes="log-container")

    # 进度显示组件调整到左侧面板下方
    with gr.Row(visible=False) as progress_row:
        gr.HTML("""
        <div class="progress-text">
            <span>当前进度：</span>
            <span id="current-step" style="color: #2b6de3;">初始化...</span>
            <span id="progress-percent" style="margin-left:15px;color: #e32b2b;">0%</span>
        </div>
        """)


    # 定义函数处理事件
    def clear_chat_history():
        return None, "对话已清空"


    def process_chat(question: str, history: Optional[List[Tuple[str, str]]], enable_web_search: bool,
                     model_choice: str):
        if history is None:
            history = []

        # 更新模型选择信息的显示
        api_text = """
        <div class="api-info" style="margin-top:10px;padding:10px;border-radius:5px;background:var(--panel-bg);border:1px solid var(--border-color);">
            <p>📢 <strong>功能说明：</strong></p>
            <p>1. <strong>联网搜索</strong>：%s</p>
            <p>2. <strong>模型选择</strong>：当前使用 <strong>%s</strong> %s</p>
        </div>
        """ % (
            "已启用" if enable_web_search else "未启用",
            "Cloud DeepSeek-R1 模型" if model_choice == "siliconflow" else "本地 Ollama 模型",
            "(需要在.env文件中配置SERPAPI_KEY)" if enable_web_search else ""
        )

        # 如果问题为空，不处理
        if not question or question.strip() == "":
            history.append(("", "问题不能为空，请输入有效问题。"))
            return history, "", api_text

        # 添加用户问题到历史
        history.append((question, ""))

        # 创建生成器
        resp_generator = stream_answer(question, enable_web_search, model_choice)

        # 流式更新回答
        for response, status in resp_generator:
            history[-1] = (question, response)
            yield history, "", api_text


    def update_api_info(enable_web_search, model_choice):
        api_text = """
        <div class="api-info" style="margin-top:10px;padding:10px;border-radius:5px;background:var(--panel-bg);border:1px solid var(--border-color);">
            <p>📢 <strong>功能说明：</strong></p>
            <p>1. <strong>联网搜索</strong>：%s</p>
            <p>2. <strong>模型选择</strong>：当前使用 <strong>%s</strong> %s</p>
        </div>
        """ % (
            "已启用" if enable_web_search else "未启用",
            "Cloud DeepSeek-R1 模型" if model_choice == "siliconflow" else "本地 Ollama 模型",
            "(需要在.env文件中配置SERPAPI_KEY)" if enable_web_search else ""
        )
        return api_text


    # 新增：真实系统监控数据获取函数
    def get_real_system_metrics():
        """获取真实的系统监控数据"""
        try:
            import psutil
            from datetime import datetime

            # CPU使用率
            cpu_percent = psutil.cpu_percent(interval=1)
            cpu_count = psutil.cpu_count(logical=False)  # 物理核心数

            # 内存使用
            mem = psutil.virtual_memory()
            memory_total = round(mem.total / (1024 ** 3), 1)  # GB
            memory_used = round(mem.used / (1024 ** 3), 1)
            memory_percent = mem.percent

            # 磁盘使用
            disk = psutil.disk_usage('/')
            disk_total = round(disk.total / (1024 ** 3), 1)
            disk_used = round(disk.used / (1024 ** 3), 1)
            disk_percent = disk.percent

            # 网络流量
            net_io = psutil.net_io_counters()
            network_up = round(net_io.bytes_sent / 1024, 1)  # KB
            network_down = round(net_io.bytes_recv / 1024, 1)

            # 向量数据库真实信息
            doc_count = len(faiss_contents_map) if faiss_contents_map else 0
            vector_count = faiss_index.ntotal if faiss_index else 0

            # 模型连接状态
            model_status = "离线"
            model_connection = "未连接"

            current_model_choice = "siliconflow"  # 默认值

            # 修复：根据当前实际使用的模型来检测状态
            def check_model_status(model_choice):
                try:
                    # 直接测试当前实际在使用的API
                    if model_choice == "siliconflow":
                        # 更健壮的SiliconFlow API测试
                        test_response = call_siliconflow_api("测试连接", max_tokens=5)
                        if test_response and isinstance(test_response, str):
                            return "在线", "SiliconFlow API正常"
                        return "在线", "API响应正常"
                    else:
                        # 只有当选择ollama时才检测本地服务
                        try:
                            response = requests.get(
                                "http://localhost:11434/api/tags",
                                timeout=10,
                                headers={'Connection': 'close'}
                            )
                            if response.status_code == 200:
                                return "在线", "Ollama服务正常"
                            return "在线", f"HTTP {response.status_code}"
                        except requests.exceptions.ConnectionError:
                            return "离线", "Ollama服务未启动"
                        except Exception as e:
                            return "在线", f"连接异常: {str(e)}"
                except Exception as e:
                    return "离线", f"检测失败: {str(e)}"

            model_status, model_connection = check_model_status(current_model_choice)

            # 生成进度条HTML
            def create_progress_bar(percent, color="var(--tech-cyan)"):
                return f"""
                <div class="progress-container">
                    <div class="progress-bar" style="width: {percent}%; background: {color}"></div>
                </div>
                """

            cpu_color = "#4CAF50" if cpu_percent < 50 else "#FFC107" if cpu_percent < 80 else "#f44336"
            cpu_progress = create_progress_bar(cpu_percent, cpu_color)

            mem_color = "#4CAF50" if memory_percent < 50 else "#FFC107" if memory_percent < 80 else "#f44336"
            memory_progress = create_progress_bar(memory_percent, mem_color)

            disk_color = "#4CAF50" if disk_percent < 50 else "#FFC107" if disk_percent < 80 else "#f44336"
            disk_progress = create_progress_bar(disk_percent, disk_color)

            # 生成日志
            log_entries = []
            current_time = datetime.now().strftime("%H:%M:%S")

            # 系统日志
            log_entries.append(f"""
            <div class="log-entry">
                <span class="log-time">[{current_time}]</span>
                <span class="log-info">[INFO]</span> 系统监控数据已更新
            </div>
            """)

            # 警告日志
            if cpu_percent > 80:
                log_entries.append(f"""
                <div class="log-entry">
                    <span class="log-time">[{current_time}]</span>
                    <span class="log-warning">[WARNING]</span> CPU使用率过高: {cpu_percent}%
                </div>
                """)

            if memory_percent > 80:
                log_entries.append(f"""
                <div class="log-entry">
                    <span class="log-time">[{current_time}]</span>
                    <span class="log-warning">[WARNING]</span> 内存使用率过高: {memory_percent}%
                </div>
                """)

            if disk_percent > 90:
                log_entries.append(f"""
                <div class="log-entry">
                    <span class="log-time">[{current_time}]</span>
                    <span class="log-error">[ERROR]</span> 磁盘空间不足: {disk_percent}%
                </div>
                """)

            log_html = "".join(log_entries[-10:])  # 只显示最近10条日志

            return (
                f"{cpu_percent}%",  # cpu_value
                cpu_progress,  # cpu_progress
                f"物理核心: {cpu_count}",  # cpu_info
                f"{memory_used}GB / {memory_total}GB",  # memory_value
                memory_progress,  # memory_progress
                f"使用率: {memory_percent}%",  # memory_info
                f"{disk_used}GB / {disk_total}GB",  # disk_value
                disk_progress,  # disk_progress
                f"使用率: {disk_percent}%",  # disk_info
                f"↑ {network_up}KB ↓ {network_down}KB",  # network_value
                f"累计流量",  # network_info
                f"{int(time.time() - psutil.boot_time())}s",  # latency_value (系统运行时间)
                f"系统运行时间",  # latency_info
                f"{doc_count + vector_count}",  # request_value
                f"文档: {doc_count} | 向量: {vector_count}",  # request_info
                f"分块数: {doc_count}",  # vector_db_value
                f"向量数: {vector_count}",  # vector_db_info
                f"状态: {model_status}",  # model_status_value
                f"连接: {model_connection}",  # model_status_info
                log_html  # log_display
            )

        except Exception as e:
            error_msg = f"监控数据获取失败: {str(e)}"
            return (
                "错误", "", error_msg,
                "错误", "", error_msg,
                "错误", "", error_msg,
                "错误", error_msg,
                "错误", error_msg,
                "错误", error_msg,
                "错误", error_msg,
                "错误", error_msg,
                f"<div class='log-error'>[ERROR] {error_msg}</div>"
            )


    # 新增：清空日志函数
    def clear_system_logs():
        return "<div class='log-info'>日志已清空</div>"


    # 绑定UI事件
    upload_btn.click(
        process_multiple_pdfs,
        inputs=[file_input],
        outputs=[upload_status, file_list],
        show_progress=True
    )

    # 绑定提问按钮
    ask_btn.click(
        process_chat,
        inputs=[question_input, chatbot, web_search_checkbox, model_choice],
        outputs=[chatbot, question_input, api_info]
    )

    # 绑定清空按钮
    clear_btn.click(
        clear_chat_history,
        inputs=[],
        outputs=[chatbot, status_display]
    )

    # 当切换联网搜索或模型选择时更新API信息
    web_search_checkbox.change(
        update_api_info,
        inputs=[web_search_checkbox, model_choice],
        outputs=[api_info]
    )

    model_choice.change(
        update_api_info,
        inputs=[web_search_checkbox, model_choice],
        outputs=[api_info]
    )

    # 新增：分块可视化刷新按钮事件
    refresh_chunks_btn.click(
        fn=get_document_chunks,
        outputs=[chunks_data, chunks_status]
    )

    # 新增：分块表格点击事件
    chunks_data.select(
        fn=show_chunk_details,
        outputs=chunk_detail_text
    )

    # 新增：系统监控刷新按钮事件 - 使用真实数据
    refresh_monitor_btn.click(
        fn=get_real_system_metrics,
        outputs=[
            cpu_value, cpu_progress, cpu_info,
            memory_value, memory_progress, memory_info,
            disk_value, disk_progress, disk_info,
            network_value, network_info,
            latency_value, latency_info,
            request_value, request_info,
            vector_db_value, vector_db_info,
            model_status_value, model_status_info,
            log_display
        ]
    )

    # 新增：清空日志按钮事件
    clear_logs_btn.click(
        fn=clear_system_logs,
        outputs=[log_display]
    )

# 修改JavaScript注入部分
demo._js = """
function gradioApp() {
    // 设置默认主题为暗色
    document.documentElement.setAttribute('data-theme', 'dark');

    const observer = new MutationObserver((mutations) => {
        document.getElementById("loading").style.display = "none";
        const progress = document.querySelector('.progress-text');
        if (progress) {
            const percent = document.querySelector('.progress > div')?.innerText || '';
            const step = document.querySelector('.progress-description')?.innerText || '';
            document.getElementById('current-step').innerText = step;
            document.getElementById('progress-percent').innerText = percent;
        }
    });
    observer.observe(document.body, {childList: true, subtree: true});
}

function toggleTheme() {
    const root = document.documentElement;
    const currentTheme = root.getAttribute('data-theme');
    const newTheme = currentTheme === 'light' ? 'dark' : 'light';
    root.setAttribute('data-theme', newTheme);
}

// 初始化主题和自动刷新
document.addEventListener('DOMContentLoaded', () => {
    document.documentElement.setAttribute('data-theme', 'dark');

    // 添加动画效果
    setTimeout(() => {
        const elements = document.querySelectorAll('.chat-message, .metric-card, .model-card');
        elements.forEach(el => {
            el.style.opacity = '0';
            el.style.transform = 'translateY(20px)';
            el.style.transition = 'opacity 0.5s ease, transform 0.5s ease';
        });

        setTimeout(() => {
            elements.forEach((el, index) => {
                setTimeout(() => {
                    el.style.opacity = '1';
                    el.style.transform = 'translateY(0)';
                }, index * 100);
            });
        }, 300);
    }, 500);

    // 系统监控页面自动刷新
    let refreshInterval;
    const monitorTab = document.querySelector('[data-testid="tab-📈 系统监控"]');
    if (monitorTab) {
        monitorTab.addEventListener('click', () => {
            // 清除现有定时器
            if (refreshInterval) clearInterval(refreshInterval);

            // 每10秒自动刷新监控数据
            refreshInterval = setInterval(() => {
                const refreshBtn = document.querySelector('button[value="🔄 刷新数据"]');
                if (refreshBtn) refreshBtn.click();
            }, 10000);
        });
    }
});
"""


# 新增：向量数据库信息获取函数（需要在vector_store模块中实现）
def get_vector_store_info():
    """获取向量数据库的真实统计信息"""
    try:
        if faiss_index is None:
            return {'document_count': 0, 'chunk_count': 0}

        # 文档数 = 唯一文档ID数量（根据元数据中的doc_id）
        doc_ids = set()
        chunk_count = 0

        # 遍历所有元数据记录
        for meta in faiss_metadatas_map.values():
            if 'doc_id' in meta:
                doc_ids.add(meta['doc_id'])
            chunk_count += 1

        # 如果元数据中没有doc_id，则使用原始ID计数
        if not doc_ids:
            doc_ids = set(faiss_contents_map.keys())

        return {
            'document_count': len(doc_ids),
            'chunk_count': faiss_index.ntotal if faiss_index else chunk_count
        }
    except Exception as e:
        logging.error(f"获取向量数据库信息失败: {str(e)}")
        return {'document_count': 0, 'chunk_count': 0}


# 修改端口检查函数
def is_port_available(port):
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.settimeout(1)
        return s.connect_ex(('127.0.0.1', port)) != 0  # 更可靠的检测方式


def check_environment():
    """环境依赖检查（云端API版本）"""
    # 检查 SiliconFlow API 密钥
    if not SILICONFLOW_API_KEY:
        print("❌ 未配置 SiliconFlow API 密钥")
        print("请在 .env 文件中设置 SILICONFLOW_API_KEY")
        return False

    print("✅ SiliconFlow API 密钥已配置")
    print("✅ 跳过本地 Ollama 检查，使用云端 API 模式")

    # 测试 SiliconFlow API 连接
    try:
        test_prompt = "你好，请回复'连接成功'"
        result = call_siliconflow_api(test_prompt, temperature=0.1, max_tokens=50)
        if "连接成功" in result or "你好" in result:
            print("✅ SiliconFlow API 连接测试成功")
            return True
        else:
            print("⚠️ SiliconFlow API 响应异常，但继续运行")
            return True
    except Exception as e:
        print(f"⚠️ SiliconFlow API 测试失败: {e}")
        print("⚠️ 继续运行，请确保 API 密钥正确")
        return True


# 方案2：禁用浏览器缓存（添加meta标签）
gr.HTML("""
<meta http-equiv="Cache-Control" content="no-cache, no-store, must-revalidate">
<meta http-equiv="Pragma" content="no-cache">
<meta http-equiv="Expires" content="0">
""")

if __name__ == "__main__":
    # 检查远程 API 和模型
    if not check_environment():
        exit(1)

    # 本地端口选择逻辑
    ports = [17995, 17996, 17997, 17998, 17999]
    selected_port = next((p for p in ports if is_port_available(p)), None)

    if not selected_port:
        print("所有端口都被占用，请手动释放端口")
        exit(1)

    try:
        # 打开浏览器
        webbrowser.open(f"http://127.0.0.1:{selected_port}")

        # 启动 Demo
        demo.launch(
            server_port=selected_port,
            server_name="0.0.0.0",
            show_error=True,
            ssl_verify=False,
            height=900
        )
    except Exception as e:
        print(f"启动失败: {str(e)}")

