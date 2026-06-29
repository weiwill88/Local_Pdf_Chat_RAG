"""
Pemuat Dokumen —— Ekstraksi Teks Dokumen Multi-Format

Poin Pembelajaran:
- Memahami metode parsing untuk various format dokumen (PDF, Word, Excel, PPT)
- Memahami langkah pertama RAG: mengonversi dokumen tidak terstruktur menjadi teks biasa
"""

import os
import logging


def extract_text(filepath):
    """
    Mengekstrak konten teks biasa dari file

    Format yang didukung: PDF / Word / Excel / PPT / Teks Biasa / Markdown

    Args:
        filepath: Jalur file

    Returns:
        String konten teks yang diekstrak
    """
    file_ext = os.path.splitext(filepath)[1].lower()

    if file_ext == '.pdf':
        try:
            from docling.document_converter import DocumentConverter, PdfFormatOption
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from config import IMAGE_PLACEHOLDER_THRESHOLD, MIN_VISIBLE_RATIO
            
            # Konversi awal tanpa OCR (cepat)
            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = False
            converter = DocumentConverter(
                format_options={
                    "pdf": PdfFormatOption(pipeline_options=pipeline_options)
                }
            )
            result = converter.convert(filepath)
            
            final_md_blocks = []
            converter_ocr = None
            
            num_pages = len(result.document.pages)
            for page_no in range(1, num_pages + 1):
                page_md = result.document.export_to_markdown(page_no=page_no)
                
                total_chars = len(page_md)
                placeholder_count = page_md.count("<!-- image -->")
                
                visible_chars = total_chars - (placeholder_count * len("<!-- image -->"))
                visible_ratio = visible_chars / max(1, total_chars)
                
                # Deteksi halaman dominan gambar (Lewati Cover page > 1)
                if placeholder_count >= IMAGE_PLACEHOLDER_THRESHOLD and visible_ratio < MIN_VISIBLE_RATIO and page_no > 1:
                    if converter_ocr is None:
                        pipeline_options_ocr = PdfPipelineOptions()
                        pipeline_options_ocr.do_ocr = True
                        converter_ocr = DocumentConverter(
                            format_options={"pdf": PdfFormatOption(pipeline_options=pipeline_options_ocr)}
                        )
                    
                    try:
                        import pypdf
                        import tempfile
                        # Ekstrak halaman tunggal untuk OCR agar lebih cepat
                        with open(filepath, "rb") as f_in:
                            reader = pypdf.PdfReader(f_in)
                            writer = pypdf.PdfWriter()
                            writer.add_page(reader.pages[page_no - 1])
                            
                            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f_out:
                                writer.write(f_out)
                                temp_pdf_path = f_out.name
                                
                        page_result = converter_ocr.convert(temp_pdf_path)
                        page_ocr_md = page_result.document.export_to_markdown()
                        os.remove(temp_pdf_path)
                        
                        new_visible_chars = len(page_ocr_md) - (page_ocr_md.count("<!-- image -->") * len("<!-- image -->"))
                        logging.info(f"Page {page_no} OCR triggered. Visible chars: {visible_chars} -> {new_visible_chars}")
                        page_md = page_ocr_md
                        
                    except ImportError:
                        logging.warning("pypdf tidak terinstal, OCR halaman penuh sebagai fallback (sangat lambat). Silakan `pip install pypdf`.")
                        # Jika pypdf tidak ada, panggil fallback dengan convert ulang (mungkin memproses seluruh file)
                        try:
                            # Coba opsi page_range docling jika didukung
                            page_result = converter_ocr.convert(filepath, page_range=(page_no, page_no))
                            page_md = page_result.document.export_to_markdown()
                        except Exception:
                            page_result = converter_ocr.convert(filepath)
                            page_md = page_result.document.export_to_markdown(page_no=page_no)
                            
                        new_visible_chars = len(page_md) - (page_md.count("<!-- image -->") * len("<!-- image -->"))
                        logging.info(f"Page {page_no} OCR triggered (Fallback). Visible chars: {visible_chars} -> {new_visible_chars}")

                final_md_blocks.append(page_md)
                
            return "\n\n".join(final_md_blocks)
        except Exception as e:
            logging.error(f"Gagal memproses dokumen PDF menggunakan docling: {e}")
            return ""

    elif file_ext in ['.txt', '.md']:
        with open(filepath, 'r', encoding='utf-8') as file:
            return file.read()

    elif file_ext == '.docx':
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        except ImportError:
            logging.error("Memproses dokumen Word memerlukan instalasi pustaka python-docx")
            return ""

    elif file_ext in ['.xlsx', '.xls']:
        try:
            import pandas as pd
            text = ""
            xl = pd.ExcelFile(filepath)
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                text += f"Lembar Kerja: {sheet_name}\n"
                text += df.to_string(index=False) + "\n\n"
            return text
        except ImportError:
            logging.error("Memproses file Excel memerlukan instalasi pustaka pandas")
            return ""

    elif file_ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            logging.error("Memproses file PPT memerlukan instalasi pustaka python-pptx")
            return ""

    else:
        logging.warning(f"Format file tidak didukung: {file_ext}")
        return ""
